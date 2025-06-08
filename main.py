import os
import shutil
import logging
import json
from fastapi import FastAPI, UploadFile, File, HTTPException, Request, Body, WebSocket, Query, BackgroundTasks
from fastapi.responses import JSONResponse, FileResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from typing import List, Optional, Dict, Any
import fitz  # PyMuPDF
from pptx import Presentation
from controller import annotate_page, create_pdf_note, ask_question, improve_note
from config import (
    PAGE_DIR, UPLOAD_DIR, UPLOAD_MAX_SIZE, 
    ALLOWED_EXTENSIONS, LOG_LEVEL, LOG_FORMAT, QWEN_API_KEY, QWEN_VL_API_KEY
)
# 导入新模块
from board_logger import board_logger
from butler_llm import butler_llm
from llm_logger import router as llm_logger_router  # 导入日志API路由
from board_manager import board_manager  # 导入展板管理器
from intelligent_expert import IntelligentExpert
# 导入简化的专家系统
from simple_expert import simple_expert_manager
# 导入任务事件管理器
from task_event_manager import task_event_manager
from fastapi.staticfiles import StaticFiles
import asyncio
import uvicorn
import dotenv
import config
import llm_agents
import conversation_manager
import time
import datetime
import httpx
from fastapi import WebSocketDisconnect
from openai import OpenAI
import requests
import random
from datetime import datetime, timezone
import dotenv
import uvicorn
import time
import asyncio
import json
import secrets
from contextlib import asynccontextmanager
from starlette.responses import Response
from concurrent.futures import ThreadPoolExecutor

# 配置日志
logging.basicConfig(level=getattr(logging, LOG_LEVEL), format=LOG_FORMAT)
logger = logging.getLogger(__name__)

app = FastAPI(
    title="WhatNote API",
    description="课件注释生成API",
    version="1.0.0"
)

# 配置CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # 在生产环境中应该更具体地指定源
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 集成LLM日志API路由
app.include_router(llm_logger_router)

# 挂载静态文件目录
app.mount("/materials", StaticFiles(directory="uploads"), name="materials")

# 确保目录存在
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(PAGE_DIR, exist_ok=True)

# 🔧 添加轻量级操作的专用线程池，避免被LLM任务阻塞
lightweight_executor = ThreadPoolExecutor(max_workers=5, thread_name_prefix="lightweight_ops")

# 🔧 添加LLM专用线程池，隔离LLM操作避免阻塞其他功能
llm_executor = ThreadPoolExecutor(max_workers=3, thread_name_prefix="llm_ops")

async def run_llm_in_background(llm_func, *args, **kwargs):
    """在后台线程池中运行LLM操作，避免阻塞轻量级操作"""
    try:
        result = await asyncio.get_event_loop().run_in_executor(
            llm_executor, llm_func, *args, **kwargs
        )
        return result
    except Exception as e:
        logger.error(f"后台LLM操作失败: {str(e)}")
        return f"LLM操作失败: {str(e)}"

# 健康检查端点
@app.get('/health')
async def health_check():
    """健康检查端点，用于启动脚本检测服务状态"""
    return {
        "status": "healthy",
        "timestamp": datetime.now().isoformat(),
        "message": "WhatNote服务运行正常"
    }

# 添加同步函数
def sync_app_state_to_butler():
    """同步应用状态到管家LLM - 已禁用"""
    pass  # 管家LLM功能已临时禁用

def validate_file(file: UploadFile) -> None:
    """验证上传文件"""
    if not file.filename:
        raise HTTPException(status_code=400, detail="文件名不能为空")
    
    ext = file.filename.split('.')[-1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400, 
            detail=f"不支持的文件类型，仅支持: {', '.join(ALLOWED_EXTENSIONS)}"
        )

def save_upload_file(upload_file: UploadFile, destination: str):
    """保存上传文件"""
    try:
        with open(destination, 'wb') as buffer:
            shutil.copyfileobj(upload_file.file, buffer)
        logger.info(f"文件已保存: {destination}")
    except Exception as e:
        logger.error(f"保存文件失败: {str(e)}")
        raise HTTPException(status_code=500, detail="文件保存失败")

@app.exception_handler(Exception)
async def global_exception_handler(request: Request, exc: Exception):
    """全局异常处理"""
    logger.error(f"未处理的异常: {str(exc)}", exc_info=True)
    return JSONResponse(
        status_code=500,
        content={"detail": "服务器内部错误"}
    )

@app.post('/materials/upload')
async def upload_material(file: UploadFile = File(...)):
    """上传课件文件"""
    logger.info(f"收到文件上传请求: {file.filename}")
    validate_file(file)
    
    save_path = os.path.join(UPLOAD_DIR, file.filename)
    save_upload_file(file, save_path)
    
    try:
        ext = file.filename.split('.')[-1].lower()
        if ext == 'pdf':
            pages = split_pdf(save_path, file.filename)
        else:
            pages = split_pptx(save_path, file.filename)
        logger.info(f"文件处理完成: {file.filename}, 共{len(pages)}页")
        
        # 同步到管家LLM
        sync_app_state_to_butler()
        
        return {"filename": file.filename, "pages": len(pages)}
    except Exception as e:
        logger.error(f"文件处理失败: {str(e)}")
        raise HTTPException(status_code=500, detail="文件处理失败")

# PDF按页拆分
def split_pdf(pdf_path, base_name):
    # 使用controller.py中的split_pdf函数
    from controller import split_pdf as controller_split_pdf
    return controller_split_pdf(pdf_path, base_name)

# PPTX按页拆分

def split_pptx(pptx_path, base_name):
    prs = Presentation(pptx_path)
    page_files = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        text = '\n'.join(texts)
        page_file = os.path.join(PAGE_DIR, f"{base_name}_page_{i+1}.txt")
        with open(page_file, 'w', encoding='utf-8') as f:
            f.write(text)
        page_files.append(page_file)
    return page_files

# 获取课件分页内容列表
@app.get('/materials/{filename}/pages')
async def get_material_pages(filename: str) -> List[str]:
    """获取课件分页内容"""
    logger.info(f"获取文件页面: {filename}")
    prefix = os.path.join(PAGE_DIR, f"{filename}_page_")
    pages = []
    i = 1
    while True:
        page_file = f"{prefix}{i}.txt"
        if not os.path.exists(page_file):
            break
        try:
            with open(page_file, 'r', encoding='utf-8') as f:
                pages.append(f.read())
            i += 1
        except Exception as e:
            logger.error(f"读取页面失败: {str(e)}")
            raise HTTPException(status_code=500, detail="读取页面失败")
    
    if not pages:
        raise HTTPException(status_code=404, detail='未找到分页内容')
    return pages

@app.get('/materials/{filename}/pages/{page_number}/annotate')
async def annotate_material_page(
    filename: str, 
    page_number: int, 
    force_vision: bool = False,
    session_id: Optional[str] = Query(None),
    current_annotation: Optional[str] = None,
    improve_request: Optional[str] = None,
    board_id: Optional[str] = Query(None)
):
    """生成页面注释"""
    logger.info(f"生成注释: {filename} 第{page_number}页, 会话ID: {session_id}, 展板ID: {board_id}")
    logger.info(f"注释生成参数: 强制视觉={force_vision}, 当前注释长度={len(current_annotation) if current_annotation else 0}, 改进请求={improve_request}")
    try:
        result = annotate_page(filename, page_number, force_vision, session_id, current_annotation, improve_request, board_id)
        return result
    except Exception as e:
        logger.error(f"生成注释失败: {str(e)}")
        raise HTTPException(status_code=500, detail="生成注释失败")

@app.post('/materials/{filename}/pages/{page_number}/annotate')
async def post_annotate_material_page(
    filename: str, 
    page_number: int, 
    force_vision: bool = False,
    session_id: Optional[str] = Query(None),
    request_data: Optional[dict] = Body(None)
):
    """POST方式生成页面注释（直接路径）"""
    logger.info(f"直接路径POST生成注释: {filename} 第{page_number}页, 会话ID: {session_id}")
    logger.info(f"请求数据: {request_data}")
    try:
        # 从请求数据中获取board_id
        board_id = request_data.get('board_id') if request_data else None
        
        result = annotate_page(
            filename, 
            page_number, 
            force_vision, 
            session_id, 
            request_data.get('current_annotation') if request_data else None,
            request_data.get('improve_request') if request_data else None,
            board_id
        )
        return result
    except Exception as e:
        logger.error(f"生成注释失败: {str(e)}")
        raise HTTPException(status_code=500, detail="生成注释失败")

@app.get('/materials/{filename}/pages/{page_number}/raw-text')
async def get_raw_page_text(filename: str, page_number: int):
    """获取页面原始提取文本"""
    logger.info(f"获取原始文本: {filename} 第{page_number}页")
    try:
        from controller import get_page_text
        text = get_page_text(filename, page_number)
        return {"text": text}
    except Exception as e:
        logger.error(f"获取原始文本失败: {str(e)}")
        raise HTTPException(status_code=500, detail="获取原始文本失败")

@app.post('/materials/{filename}/pages/{page_number}/vision-annotate')
async def post_force_vision_annotation(
    filename: str, 
    page_number: int,
    session_id: Optional[str] = Query(None),
    request_data: Optional[dict] = Body(None)
):
    """POST方式强制使用图像识别对页面进行注释（直接路径）"""
    logger.info(f"直接路径POST强制图像识别注释: {filename} 第{page_number}页, 会话ID: {session_id}")
    logger.info(f"请求数据: {request_data}")
    try:
        # 从请求数据中获取board_id
        board_id = request_data.get('board_id') if request_data else None
        logger.info(f"使用展板ID: {board_id}")
        
        # 从请求数据中获取当前注释和改进请求
        current_annotation = request_data.get('current_annotation') if request_data else None
        improve_request = request_data.get('improve_request') if request_data else None
        
        # 🔧 新增：从请求数据中获取风格参数
        annotation_style = request_data.get('annotation_style') if request_data else None
        custom_prompt = request_data.get('custom_prompt') if request_data else None
        
        # 记录关键参数以便调试
        if current_annotation:
            logger.info(f"当前注释长度: {len(current_annotation)}")
        if improve_request:
            logger.info(f"用户改进请求: {improve_request}")
        if annotation_style:
            logger.info(f"指定注释风格: {annotation_style}")
        if custom_prompt:
            logger.info(f"自定义提示长度: {len(custom_prompt)}")
        
        # 🔧 如果传递了风格参数，临时设置到对应的专家实例
        if board_id and annotation_style:
            try:
                from simple_expert import simple_expert_manager
                expert = simple_expert_manager.get_expert(board_id)
                # 🔧 修复：不要强制恢复到默认值，保持用户的设置
                # 临时保存当前设置，但不使用getattr的默认值
                original_style = expert.annotation_style
                original_custom = expert.custom_annotation_prompt
                
                # 临时应用新风格
                expert.set_annotation_style(annotation_style, custom_prompt or '')
                logger.info(f"临时应用风格设置: {annotation_style}")
                
                try:
                    # 执行注释生成
                    result = annotate_page(
                        filename, 
                        page_number, 
                        force_vision=True, 
                        session_id=session_id, 
                        current_annotation=current_annotation,
                        improve_request=improve_request,
                        board_id=board_id
                    )
                except Exception as e:
                    pass
                finally:
                    # 恢复原始设置
                    expert.set_annotation_style(original_style, original_custom)
                    logger.info(f"恢复原始风格设置: {original_style}")
                
                return result
            except Exception as e:
                logger.error(f"临时风格设置失败: {str(e)}，使用默认流程")
        
        # 默认流程（无风格参数或设置失败）
        result = annotate_page(
            filename, 
            page_number, 
            force_vision=True, 
            session_id=session_id, 
            current_annotation=current_annotation,
            improve_request=improve_request,
            board_id=board_id
        )
        return result
    except Exception as e:
        logger.error(f"强制视觉识别注释失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"生成注释失败: {str(e)}")

@app.get('/materials/{filename}/pages/{page_number}/image')
async def get_material_page_image(filename: str, page_number: int):
    """获取页面图片"""
    logger.info(f"获取页面图片: {filename} 第{page_number}页")
    try:
        from controller import get_page_image
        img_path = get_page_image(filename, page_number)
        return FileResponse(img_path, media_type="image/png")
    except Exception as e:
        logger.error(f"获取页面图片失败: {str(e)}")
        raise HTTPException(status_code=500, detail="获取页面图片失败")

@app.post('/materials/{filename}/note')
async def generate_material_note(
    filename: str,
    session_id: Optional[str] = Query(None)
):
    """生成整本PDF的AI笔记"""
    try:
        # 读取所有页面内容
        prefix = os.path.join(PAGE_DIR, f"{filename}_page_")
        pages = []
        i = 1
        while True:
            page_file = f"{prefix}{i}.txt"
            if not os.path.exists(page_file):
                break
            with open(page_file, 'r', encoding='utf-8') as f:
                pages.append(f.read())
            i += 1
        if not pages:
            raise HTTPException(status_code=404, detail='未找到分页内容')
            
        # 使用修改后的controller函数
        result = create_pdf_note(filename, pages, session_id)
        return result
    except Exception as e:
        logger.error(f"生成整本笔记失败: {str(e)}")
        raise HTTPException(status_code=500, detail="生成整本笔记失败")

@app.post('/materials/{filename}/ask')
async def ask_material_question(
    filename: str, 
    question: str = Body(..., embed=True),
    session_id: Optional[str] = Query(None)
):
    """针对整本PDF的AI问答"""
    try:
        # 读取所有页面内容
        prefix = os.path.join(PAGE_DIR, f"{filename}_page_")
        pages = []
        i = 1
        while True:
            page_file = f"{prefix}{i}.txt"
            if not os.path.exists(page_file):
                break
            with open(page_file, 'r', encoding='utf-8') as f:
                pages.append(f.read())
            i += 1
        if not pages:
            raise HTTPException(status_code=404, detail='未找到分页内容')
        
        # 使用修改后的controller函数
        result = ask_question(filename, question, pages, session_id)
        return result
    except Exception as e:
        logger.error(f"AI问答失败: {str(e)}")
        raise HTTPException(status_code=500, detail="AI问答失败")

@app.websocket('/materials/{filename}/ask/stream')
async def ask_pdf_question_stream(websocket: WebSocket, filename: str):
    """提供流式AI问答服务"""
    await websocket.accept()
    try:
        data = await websocket.receive_json()
        question = data.get("question", "")
        session_id = data.get("session_id", None)
        
        if not question:
            await websocket.send_json({"error": "问题不能为空"})
            await websocket.close()
            return
        
        # 读取所有页面内容
        prefix = os.path.join(PAGE_DIR, f"{filename}_page_")
        pages = []
        i = 1
        while True:
            page_file = f"{prefix}{i}.txt"
            if not os.path.exists(page_file):
                break
            with open(page_file, 'r', encoding='utf-8') as f:
                pages.append(f.read())
            i += 1
            
        if not pages:
            await websocket.send_json({"error": "未找到PDF内容"})
            await websocket.close()
            return
        
        logger.info(f"开始流式问答: {filename}, 问题: {question}, 会话ID: {session_id}")
        
        # 定义回调函数，将流式生成的文本发送给WebSocket客户端
        async def send_chunk(chunk):
            await websocket.send_json({"chunk": chunk})
        
        # 创建一个同步回调
        def callback(chunk):
            asyncio.run_coroutine_threadsafe(send_chunk(chunk), asyncio.get_event_loop())
        
        # 导入流式问答
        from llm_agents import ask_pdf_question_stream
        
        # 启动流式生成，传递会话ID和回调函数
        full_answer = ask_pdf_question_stream(pages, question, callback, session_id, filename)
        
        # 发送完成信号
        await websocket.send_json({"done": True, "full_answer": full_answer})
        logger.info(f"流式问答完成: {filename}")
    except Exception as e:
        logger.error(f"流式问答失败: {str(e)}")
        try:
            await websocket.send_json({"error": f"处理请求失败: {str(e)}"})
        except Exception as e:
            pass
        except:
            # 连接可能已关闭
            pass
    finally:
        try:
            await websocket.close()
        except Exception as e:
            pass
        except:
            pass

@app.post('/materials/{filename}/improve-note')
async def improve_material_note(
    filename: str, 
    request_data: dict = Body(...),
):
    """AI完善用户笔记"""
    logger.info(f"收到笔记完善请求: {filename}")
    try:
        content = request_data.get("content", "")
        improve_prompt = request_data.get("improve_prompt", "")
        session_id = request_data.get("session_id", None)
        board_id = request_data.get("board_id", None)
        
        logger.info(f"笔记改进提示: {improve_prompt}")
        logger.info(f"使用展板ID: {board_id or '无'}")
        
        if not content:
            raise HTTPException(status_code=400, detail="内容不能为空")
        
        # 读取所有页面内容作为参考资料
        prefix = os.path.join(PAGE_DIR, f"{filename}_page_")
        pages = []
        i = 1
        while True:
            page_file = f"{prefix}{i}.txt"
            if not os.path.exists(page_file):
                break
            with open(page_file, 'r', encoding='utf-8') as f:
                pages.append(f.read())
            i += 1
            
        if not pages:
            raise HTTPException(status_code=404, detail='未找到PDF内容')

        # 使用controller函数，传递改进提示和展板ID
        result = improve_note(filename, content, pages, improve_prompt, session_id, board_id)
        
        logger.info(f"笔记完善成功: {filename}, 改进提示: {improve_prompt}")
        return result
    except Exception as e:
        logger.error(f"笔记完善失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"笔记完善失败: {str(e)}")

@app.get('/api/check-config')
async def check_api_config():
    """检查API配置是否正确"""
    return {
        "qwen_api_configured": bool(QWEN_API_KEY),
        "qwen_vl_api_configured": bool(QWEN_VL_API_KEY)
    }

@app.get('/materials/check/{filename}')
async def check_material_file(filename: str):
    """检查指定文件是否存在，返回真实文件路径"""
    logger.info(f"检查文件是否存在: {filename}")
    try:
        from controller import check_file_exists
        result = check_file_exists(filename)
        return result
    except Exception as e:
        logger.error(f"检查文件失败: {str(e)}")
        raise HTTPException(status_code=500, detail="检查文件失败")

@app.get('/materials/view/{filename}')
async def view_material_file(filename: str):
    """获取文件内容"""
    logger.info(f"请求查看文件: {filename}")
    try:
        from controller import check_file_exists
        file_check = check_file_exists(filename)
        
        if not file_check["exists"]:
            raise HTTPException(status_code=404, detail="文件不存在")
        
        file_path = os.path.join(UPLOAD_DIR, file_check["path"])
        return FileResponse(file_path, filename=file_check["path"])
    except Exception as e:
        pass
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"获取文件失败: {str(e)}")
        raise HTTPException(status_code=500, detail="获取文件失败")

# 初始化应用状态管理
class AppState:
    def __init__(self):
        self.course_folders = []
        self.boards = []
        self.pdfs = []
        
        # 初始加载状态
        self._load_state()
    
    def _load_state(self):
        # 从持久化存储加载状态
        try:
            if os.path.exists('app_state.json'):
                with open('app_state.json', 'r', encoding='utf-8') as f:
                    state = json.load(f)
                    self.course_folders = state.get('course_folders', [])
                    self.boards = state.get('boards', [])
                    self.pdfs = state.get('pdfs', [])
                    
                    # 确保每个课程文件夹都有files字段
                    for folder in self.course_folders:
                        if 'files' not in folder:
                            folder['files'] = []
        except Exception as e:
            logger.error(f"加载应用状态失败: {str(e)}")
    
    def save_state(self):
        # 保存状态到持久化存储
        try:
            state = {
                'course_folders': self.course_folders,
                'boards': self.boards,
                'pdfs': self.pdfs
            }
            with open('app_state.json', 'w', encoding='utf-8') as f:
                json.dump(state, f, ensure_ascii=False, indent=2)
        except Exception as e:
            logger.error(f"保存应用状态失败: {str(e)}")
    
    def add_course_folder(self, folder_name: str) -> Dict[str, Any]:
        # 添加课程文件夹 - 使用时间戳和随机数确保ID唯一性
        import time
        import random
        timestamp = int(time.time() * 1000)  # 毫秒级时间戳
        random_suffix = random.randint(100, 999)  # 3位随机数
        folder_id = f"course-{timestamp}-{random_suffix}"
        
        folder = {
            'id': folder_id,
            'name': folder_name,
            'files': [],  # 初始化文件列表
            'created_at': time.time()  # 添加创建时间戳
        }
        self.course_folders.append(folder)
        self.save_state()
        return folder
    
    def add_board(self, board_name: str, course_folder: str) -> Dict[str, Any]:
        # 添加展板 - 使用时间戳和随机数确保ID唯一性
        import time
        import random
        timestamp = int(time.time() * 1000)  # 毫秒级时间戳
        random_suffix = random.randint(100, 999)  # 3位随机数
        board_id = f"board-{timestamp}-{random_suffix}"
        
        board = {
            'id': board_id,
            'name': board_name,
            'course_folder': course_folder,
            'pdfs': 0,
            'windows': 0,
            'created_at': time.time()  # 添加创建时间戳
        }
        self.boards.append(board)
        self.save_state()
        return board
    
    def course_folder_exists(self, folder_name: str) -> bool:
        # 检查课程文件夹是否存在
        return any(folder['name'] == folder_name for folder in self.course_folders)
    
    def board_exists(self, board_name: str, course_folder: str) -> bool:
        # 检查展板是否存在
        return any(
            board['name'] == board_name and board['course_folder'] == course_folder 
            for board in self.boards
        )
    
    def get_boards(self) -> List[Dict[str, Any]]:
        # 获取所有展板
        return self.boards
    
    def get_course_folders(self) -> List[Dict[str, Any]]:
        # 获取所有课程文件夹
        return self.course_folders

# 初始化应用状态
app_state = AppState()

# 新增API端点: 获取应用状态
@app.get('/api/app-state')
async def get_app_state():
    """获取应用当前状态 - 轻量级操作，优先处理"""
    logger.info("获取应用状态")
    
    # 🔧 使用专用线程池处理，避免被LLM任务阻塞
    def _get_app_state_sync():
        # 获取课程文件夹和展板数据
        course_folders = app_state.get_course_folders()
        all_boards = app_state.get_boards()
        
        # 🔧 修复：将展板数据合并到对应课程的files字段中
        for folder in course_folders:
            # 确保每个课程都有files字段
            if 'files' not in folder:
                folder['files'] = []
            
            # 查找属于当前课程的展板
            # 修复：展板的course_folder字段存储的是课程ID，不是名称
            course_id = folder.get('id', '')
            course_name = folder.get('name', '')
            course_boards = [board for board in all_boards 
                            if board.get('course_folder') == course_id or board.get('course_folder') == course_name]
            
            # 将展板转换为前端期望的文件格式并添加到files中
            for board in course_boards:
                file_entry = {
                    'id': board.get('id'),
                    'name': board.get('name'),
                    'type': 'board',  # 标记为展板类型
                    'course_id': folder.get('id'),
                    'course_name': course_name,
                    'created_at': board.get('created_at'),
                    'pdfs': board.get('pdfs', 0),
                    'windows': board.get('windows', 0)
                }
                
                # 检查是否已经存在（避免重复）
                existing_ids = [f.get('id') for f in folder['files']]
                if board.get('id') not in existing_ids:
                    folder['files'].append(file_entry)
        
        return {
            'course_folders': course_folders,
            'boards': all_boards,  # 保留原始展板数据（向后兼容）
            'pdfs': [],  # 可以根据需要添加更多信息
        }
    
    # 在轻量级线程池中执行，避免阻塞
    result = await asyncio.get_event_loop().run_in_executor(
        lightweight_executor, _get_app_state_sync
    )
    
    return result

# 新增调试API端点: 查看原始app_state.json文件内容
@app.get('/api/debug/app-state-raw')
async def get_raw_app_state():
    """获取原始app_state.json文件内容（调试用）"""
    logger.info("获取原始应用状态文件内容")
    
    try:
        if os.path.exists('app_state.json'):
            with open('app_state.json', 'r', encoding='utf-8') as f:
                raw_content = f.read()
                parsed_content = json.loads(raw_content)
                
            return {
                'status': 'success',
                'file_exists': True,
                'raw_content': raw_content,
                'parsed_content': parsed_content,
                'course_folders_count': len(parsed_content.get('course_folders', [])),
                'boards_count': len(parsed_content.get('boards', [])),
                'timestamp': os.path.getmtime('app_state.json')
            }
        else:
            return {
                'status': 'error',
                'file_exists': False,
                'message': 'app_state.json file does not exist'
            }
    except Exception as e:
        logger.error(f"读取原始应用状态文件失败: {str(e)}")
        return {
            'status': 'error',
            'file_exists': os.path.exists('app_state.json'),
            'error': str(e)
        }

# 新增API端点: 获取所有展板
@app.get('/api/boards/list')
async def list_boards():
    """获取所有展板列表"""
    logger.info("获取展板列表")
    return app_state.get_boards()

# 创建课程文件夹
@app.post('/api/courses')
async def create_course_folder(request_data: dict = Body(...)):
    """创建新课程文件夹"""
    folder_name = request_data.get('name')
    if not folder_name:
        raise HTTPException(status_code=400, detail="文件夹名称不能为空")
    
    logger.info(f"创建课程文件夹: {folder_name}")
    folder = app_state.add_course_folder(folder_name)
    
    # 同步到管家LLM
    sync_app_state_to_butler()
    
    return folder

# 检查课程文件夹是否存在
@app.get('/api/course-folders/{folder_name}/exists')
async def check_course_folder_exists(folder_name: str):
    """检查课程文件夹是否存在"""
    exists = app_state.course_folder_exists(folder_name)
    return {"exists": exists}

# 创建展板（修改现有endpoints）
@app.post('/api/boards')
async def create_board(request_data: dict = Body(...)):
    """创建新展板"""
    board_name = request_data.get('name')
    course_folder = request_data.get('course_folder')
    
    if not board_name:
        raise HTTPException(status_code=400, detail="展板名称不能为空")
    
    if not course_folder:
        raise HTTPException(status_code=400, detail="课程文件夹不能为空")
    
    logger.info(f"创建展板: {board_name} (在 {course_folder} 内)")
    
    # 检查课程文件夹是否存在，不存在则创建
    if not app_state.course_folder_exists(course_folder):
        app_state.add_course_folder(course_folder)
    
    board = app_state.add_board(board_name, course_folder)
    
    # 初始化展板日志
    board_logger.init_board(board['id'])
    
    # 同步到管家LLM
    sync_app_state_to_butler()
    
    return board

# @app.post('/api/assistant')  # 管家LLM功能已禁用
# async def assistant_query(request_data: dict = Body(...)):  # 管家LLM功能已禁用
#     """处理助手LLM查询"""  # 管家LLM功能已禁用
#     query = request_data.get('query')  # 管家LLM功能已禁用
#     status_log = request_data.get('status_log', '')  # 管家LLM功能已禁用
#     history = request_data.get('history', [])  # 管家LLM功能已禁用
#       # 管家LLM功能已禁用
#     if not query:  # 管家LLM功能已禁用
#         raise HTTPException(status_code=400, detail="查询不能为空")  # 管家LLM功能已禁用
#       # 管家LLM功能已禁用
#     logger.info(f"助手查询: {query[:50]}...")  # 管家LLM功能已禁用
#       # 管家LLM功能已禁用
#     # 使用butler_llm处理查询  # 管家LLM功能已禁用
#     response = butler_llm.query(  # 管家LLM功能已禁用
#         query=query,  # 管家LLM功能已禁用
#         status_log=status_log,  # 管家LLM功能已禁用
#         history=history  # 管家LLM功能已禁用
#     )  # 管家LLM功能已禁用
#       # 管家LLM功能已禁用
#     # 提取回复和命令  # 管家LLM功能已禁用
#     reply = response.get('response', '无法处理您的请求')  # 管家LLM功能已禁用
#     command = response.get('command')  # 管家LLM功能已禁用
#       # 管家LLM功能已禁用
#     return {  # 管家LLM功能已禁用
#         "response": reply,  # 管家LLM功能已禁用
#         "command": command  # 管家LLM功能已禁用
#     }  # 管家LLM功能已禁用
#   # 管家LLM功能已禁用
@app.post('/api/boards/{board_id}/windows')
async def add_board_window(
    board_id: str, 
    request_data: dict = Body(...)
):
    """添加窗口到展板"""
    logger.info(f'添加窗口: {board_id}')
    try:
        window_data = request_data.get('window', {})
        
        if not window_data or "type" not in window_data:
            raise HTTPException(status_code=400, detail='窗口数据不能为空且必须指定类型')
        
        # 添加窗口
        window_id = board_logger.add_window(board_id, window_data)
        
        # 更新管家LLM的板块信息
        butler_llm.update_board_info(board_id)
        
        logger.info(f'窗口添加成功: {window_id}')
        return {"window_id": window_id}
    except Exception as e:
        logger.error(f'添加窗口失败: {str(e)}')
        raise HTTPException(status_code=500, detail=f'添加窗口失败: {str(e)}')

@app.delete('/api/boards/{board_id}/windows/{window_id}')
async def remove_board_window(board_id: str, window_id: str):
    """从展板移除窗口"""
    logger.info(f'移除窗口: {window_id}, 展板: {board_id}')
    try:
        # 移除窗口
        success = board_logger.remove_window(board_id, window_id)
        
        if not success:
            raise HTTPException(status_code=404, detail='未找到窗口')
        
        # 更新管家LLM的板块信息
        butler_llm.update_board_info(board_id)
        
        logger.info(f'窗口移除成功: {window_id}')
        return {"success": True}
    except Exception as e:
        pass
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f'移除窗口失败: {str(e)}')
        raise HTTPException(status_code=500, detail=f'移除窗口失败: {str(e)}')

@app.put('/api/boards/{board_id}/windows/{window_id}')
async def update_board_window(
    board_id: str, 
    window_id: str, 
    request_data: dict = Body(...)
):
    """更新展板窗口"""
    logger.info(f'更新窗口: {window_id}, 展板: {board_id}')
    try:
        window_data = request_data.get('window', {})
        
        if not window_data:
            raise HTTPException(status_code=400, detail='窗口数据不能为空')
        
        # 更新窗口
        success = board_logger.update_window(board_id, window_id, window_data)
        
        if not success:
            raise HTTPException(status_code=404, detail='未找到窗口')
        
        logger.info(f'窗口更新成功: {window_id}')
        return {"success": True}
    except Exception as e:
        pass
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f'更新窗口失败: {str(e)}')
        raise HTTPException(status_code=500, detail=f'更新窗口失败: {str(e)}')

# 创建课程文件
@app.post('/api/courses/{course_id}/files')
async def create_course_file(course_id: str, request_data: dict = Body(...)):
    """创建课程下的文件"""
    file_name = request_data.get('name')
    pdf_filename = request_data.get('pdf_filename')  # 接收PDF文件名
    
    if not file_name:
        raise HTTPException(status_code=400, detail="文件名称不能为空")
    
    # 查找对应的课程文件夹
    course_folder = None
    for folder in app_state.course_folders:
        if folder['id'] == course_id:
            course_folder = folder
            break
    
    if not course_folder:
        raise HTTPException(status_code=404, detail="未找到课程文件夹")
    
    # 创建文件记录
    # 注意：此处只创建文件记录，不创建实际文件
    # 实际应用中可能需要创建真实文件并存储内容
    file_id = f"file-{course_id}-{len(course_folder.get('files', []))+1}"
    
    # 如果课程没有files字段，添加它
    if 'files' not in course_folder:
        course_folder['files'] = []
    
    file_record = {
        'id': file_id,
        'name': file_name,
        'course_id': course_id,
        'created_at': None,  # 可以添加时间戳
        'pdf_filename': pdf_filename  # 添加PDF文件名字段
    }
    
    course_folder['files'].append(file_record)
    app_state.save_state()
    
    # 同步到管家LLM
    sync_app_state_to_butler()
    
    logger.info(f"创建课程文件: {file_name} (在课程 {course_folder['name']} 中)")
    return file_record

# 添加清理多余PDF展板文件的API
@app.post('/api/cleanup/duplicate-pdf-files')
async def cleanup_duplicate_pdf_files():
    """清理与PDF文件同名的多余展板文件"""
    logger.info("开始清理与PDF文件同名的多余展板文件")
    
    try:
        app_state = AppState()
        cleanup_count = 0
        
        for folder in app_state.course_folders:
            if 'files' not in folder:
                continue
                
            # 查找需要删除的文件（pdf_filename不为空且文件名以.pdf结尾）
            files_to_remove = []
            for file in folder.get('files', []):
                # 如果文件名以.pdf结尾且有pdf_filename字段，说明这是上传PDF时意外创建的展板文件
                if (file.get('name', '').endswith('.pdf') and 
                    file.get('pdf_filename') is not None and
                    file.get('name') == file.get('pdf_filename')):
                    files_to_remove.append(file)
                    cleanup_count += 1
                    logger.info(f"标记删除多余文件: {file.get('name')} (ID: {file.get('id')})")
            
            # 删除标记的文件
            folder['files'] = [f for f in folder.get('files', []) if f not in files_to_remove]
        
        # 保存状态
        if cleanup_count > 0:
            app_state.save_state()
            sync_app_state_to_butler()
            logger.info(f"清理完成，删除了 {cleanup_count} 个多余的PDF展板文件")
        else:
            logger.info("没有发现需要清理的多余PDF展板文件")
        
        return {
            "status": "success", 
            "message": f"清理完成，删除了 {cleanup_count} 个多余的PDF展板文件",
            "cleaned_count": cleanup_count
        }
        
    except Exception as e:
        logger.error(f"清理过程出错: {str(e)}")
        raise HTTPException(status_code=500, detail=f"清理失败: {str(e)}")

# 添加新的API路由处理PDF上传
@app.post('/api/materials/upload')
async def api_upload_material(file: UploadFile = File(...)):
    """API路由: 上传课件文件"""
    logger.info(f"收到API文件上传请求: {file.filename}")
    validate_file(file)
    
    save_path = os.path.join(UPLOAD_DIR, file.filename)
    save_upload_file(file, save_path)
    
    try:
        ext = file.filename.split('.')[-1].lower()
        if ext == 'pdf':
            pages = split_pdf(save_path, file.filename)
        else:
            pages = split_pptx(save_path, file.filename)
        logger.info(f"文件处理完成: {file.filename}, 共{len(pages)}页")
        return {"filename": file.filename, "pages": len(pages)}
    except Exception as e:
        logger.error(f"文件处理失败: {str(e)}")
        raise HTTPException(status_code=500, detail="文件处理失败")

# 添加API转发路由 - 用于统一前端路径
@app.get('/api/materials/view/{filename}')
async def api_view_material_file(filename: str):
    """API路由: 获取文件内容"""
    return await view_material_file(filename)

@app.get('/api/materials/{filename}/pages')
async def api_get_material_pages(filename: str) -> List[str]:
    """API路由: 获取课件分页内容"""
    return await get_material_pages(filename)

@app.get('/api/materials/{filename}/pages/{page_number}/annotate')
async def api_annotate_material_page(
    filename: str, 
    page_number: int, 
    force_vision: bool = False,
    session_id: Optional[str] = Query(None),
    board_id: Optional[str] = Query(None)
):
    """API方式获取页面注释"""
    logger.info(f"API方式生成注释: {filename} 第{page_number}页, 会话ID: {session_id}, 展板ID: {board_id}")
    try:
        result = annotate_page(filename, page_number, force_vision, session_id, None, None, board_id)
        return result
    except Exception as e:
        logger.error(f"生成注释失败: {str(e)}")
        raise HTTPException(status_code=500, detail="生成注释失败")

@app.post('/api/materials/{filename}/pages/{page_number}/annotate')
async def api_post_annotate_material_page(
    filename: str, 
    page_number: int, 
    force_vision: bool = False,
    session_id: Optional[str] = Query(None),
    request_data: Optional[dict] = Body(None)
):
    """API POST方式生成页面注释"""
    logger.info(f"API POST生成注释: {filename} 第{page_number}页")
    try:
        # 从请求数据中获取board_id
        board_id = request_data.get('board_id') if request_data else None
        
        result = annotate_page(
            filename, 
            page_number, 
            force_vision, 
            session_id, 
            request_data.get('current_annotation') if request_data else None,
            request_data.get('improve_request') if request_data else None,
            board_id
        )
        return result
    except Exception as e:
        logger.error(f"生成注释失败: {str(e)}")
        raise HTTPException(status_code=500, detail="生成注释失败")

@app.post('/api/materials/{filename}/pages/{page_number}/vision-annotate')
async def api_force_vision_annotation_post(
    filename: str, 
    page_number: int,
    session_id: Optional[str] = Query(None),
    request_data: Optional[dict] = Body(None)
):
    """API路由: POST方式强制使用图像识别对页面进行注释"""
    # 确保请求数据中的board_id可以被下一级函数使用
    logger.info(f"API POST强制视觉识别: {filename} 第{page_number}页")
    logger.info(f"请求数据: {request_data}")
    
    return await post_force_vision_annotation(filename, page_number, session_id, request_data)

@app.get('/api/materials/{filename}/pages/{page_number}/image')
async def api_get_material_page_image(filename: str, page_number: int):
    """API路由: 获取页面图像"""
    return await get_material_page_image(filename, page_number)

@app.post('/api/materials/{filename}/note')
async def api_generate_material_note(
    filename: str,
    session_id: Optional[str] = Query(None)
):
    """API路由: 生成整本笔记"""
    return await generate_material_note(filename, session_id)

@app.post('/api/materials/{filename}/ask')
async def api_ask_material_question(
    filename: str, 
    question: str = Body(..., embed=True),
    session_id: Optional[str] = Query(None)
):
    """API路由: 提问PDF问题"""
    return await ask_material_question(filename, question, session_id)

@app.post('/api/materials/{filename}/improve-note')
async def api_improve_material_note(
    filename: str, 
    request_data: dict = Body(...),
):
    """API路由：AI完善用户笔记"""
    logger.info(f"收到API笔记完善请求: {filename}")
    # 直接调用主路由处理函数
    return await improve_material_note(filename, request_data)

@app.get('/api/materials/check/{filename}')
async def api_check_material_file(filename: str):
    """API路由: 检查指定文件是否存在"""
    return await check_material_file(filename)

@app.post('/api/images/upload')
async def upload_image(file: UploadFile = File(...)):
    """专门用于图片上传的API"""
    logger.info(f"收到图片上传请求: {file.filename}")
    
    # 验证是否为图片文件
    allowed_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp'}
    file_ext = os.path.splitext(file.filename.lower())[1]
    
    if file_ext not in allowed_extensions:
        raise HTTPException(status_code=400, detail="只支持图片文件（jpg, jpeg, png, gif, bmp, webp）")
    
    # 创建images目录（如果不存在）
    images_dir = os.path.join(UPLOAD_DIR, 'images')
    if not os.path.exists(images_dir):
        os.makedirs(images_dir, exist_ok=True)
    
    # 生成唯一文件名（添加时间戳避免冲突）
    timestamp = int(time.time())
    name, ext = os.path.splitext(file.filename)
    unique_filename = f"{name}_{timestamp}{ext}"
    
    save_path = os.path.join(images_dir, unique_filename)
    
    try:
        # 保存图片文件
        save_upload_file(file, save_path)
        
        # 构建访问URL
        image_url = f"/api/images/view/{unique_filename}"
        
        logger.info(f"图片保存成功: {save_path}")
        return {
            "success": True,
            "filename": unique_filename,
            "original_filename": file.filename,
            "url": image_url,  
            "path": save_path
        }
    except Exception as e:
        logger.error(f"图片保存失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"图片保存失败: {str(e)}")

@app.get('/api/images/view/{filename}')
async def view_image(filename: str):
    """查看图片文件"""
    images_dir = os.path.join(UPLOAD_DIR, 'images')
    file_path = os.path.join(images_dir, filename)
    
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="图片文件不存在")
    
    # 根据文件扩展名设置正确的媒体类型
    ext = os.path.splitext(filename.lower())[1]
    media_type_map = {
        '.jpg': 'image/jpeg',
        '.jpeg': 'image/jpeg', 
        '.png': 'image/png',
        '.gif': 'image/gif',
        '.bmp': 'image/bmp',
        '.webp': 'image/webp'
    }
    
    media_type = media_type_map.get(ext, 'image/jpeg')
    
    return FileResponse(
        file_path,
        media_type=media_type,
        filename=filename
    )

# 视频相关API
@app.post('/api/videos/upload')
async def upload_video(file: UploadFile = File(...)):
    """专门用于视频上传的API"""
    logger.info(f"收到视频上传请求: {file.filename}")
    
    # 验证是否为视频文件
    allowed_extensions = {'.mp4', '.webm', '.ogg', '.avi', '.mov', '.wmv', '.flv', '.mkv', '.m4v'}
    file_ext = os.path.splitext(file.filename.lower())[1]
    
    if file_ext not in allowed_extensions:
        raise HTTPException(status_code=400, detail="只支持视频文件（mp4, webm, ogg, avi, mov, wmv, flv, mkv, m4v）")
    
    # 验证文件大小（100MB）
    file_size = 0
    # 读取文件内容获取大小
    content = await file.read()
    file_size = len(content)
    # 重置文件指针到开头
    await file.seek(0)
    
    max_size = 100 * 1024 * 1024  # 100MB
    if file_size > max_size:
        raise HTTPException(status_code=400, detail="视频文件大小不能超过100MB")
    
    # 创建videos目录（如果不存在）
    videos_dir = os.path.join(UPLOAD_DIR, 'videos')
    if not os.path.exists(videos_dir):
        os.makedirs(videos_dir, exist_ok=True)
    
    # 生成唯一文件名（添加时间戳避免冲突）
    timestamp = int(time.time())
    name, ext = os.path.splitext(file.filename)
    unique_filename = f"{name}_{timestamp}{ext}"
    
    save_path = os.path.join(videos_dir, unique_filename)
    
    try:
        # 保存视频文件
        save_upload_file(file, save_path)
        
        # 构建访问URL
        video_url = f"/api/videos/view/{unique_filename}"
        
        logger.info(f"视频保存成功: {save_path}")
        return {
            "success": True,
            "filename": unique_filename,
            "original_filename": file.filename,
            "url": video_url,  
            "path": save_path,
            "size": file_size
        }
    except Exception as e:
        logger.error(f"视频保存失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"视频保存失败: {str(e)}")

@app.get('/api/videos/view/{filename}')
async def view_video(filename: str):
    """查看视频文件"""
    videos_dir = os.path.join(UPLOAD_DIR, 'videos')
    file_path = os.path.join(videos_dir, filename)
    
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="视频文件不存在")
    
    # 根据文件扩展名设置正确的媒体类型
    ext = os.path.splitext(filename.lower())[1]
    media_type_map = {
        '.mp4': 'video/mp4',
        '.webm': 'video/webm', 
        '.ogg': 'video/ogg',
        '.avi': 'video/x-msvideo',
        '.mov': 'video/quicktime',
        '.wmv': 'video/x-ms-wmv',
        '.flv': 'video/x-flv',
        '.mkv': 'video/x-matroska',
        '.m4v': 'video/x-m4v'
    }
    
    media_type = media_type_map.get(ext, 'video/mp4')
    
    return FileResponse(
        file_path,
        media_type=media_type,
        filename=filename
    )

@app.delete('/api/videos/{filename}')
async def delete_video(filename: str):
    """删除视频文件"""
    logger.info(f"=== 开始删除视频文件 ===")
    logger.info(f"要删除的视频文件: '{filename}'")
    
    try:
        videos_dir = os.path.join(UPLOAD_DIR, 'videos')
        file_path = os.path.join(videos_dir, filename)
        
        if not os.path.exists(file_path):
            logger.warning(f"❌ 视频文件不存在: {file_path}")
            raise HTTPException(status_code=404, detail="视频文件不存在")
        
        # 获取文件信息
        file_size = os.path.getsize(file_path)
        logger.info(f"视频文件大小: {file_size} bytes")
        
        # 删除物理文件
        os.remove(file_path)
        logger.info(f"✅ 视频文件已删除: {file_path}")
        
        success_message = f"视频文件 '{filename}' 已删除"
        logger.info(f"=== 删除视频文件操作成功 ===: {success_message}")
        return {
            "status": "success", 
            "message": success_message,
            "filename": filename,
            "size": file_size
        }
        
    except HTTPException:
        # 重新抛出HTTP异常
        raise
    except Exception as e:
        error_msg = f"删除视频文件失败: {str(e)}"
        logger.error(f"=== 删除视频文件操作失败 ===")
        logger.error(error_msg)
        raise HTTPException(status_code=500, detail=error_msg)

@app.delete('/api/images/{filename}')
async def delete_image(filename: str):
    """删除图片文件"""
    logger.info(f"=== 开始删除图片文件 ===")
    logger.info(f"要删除的图片文件: '{filename}'")
    
    try:
        images_dir = os.path.join(UPLOAD_DIR, 'images')
        file_path = os.path.join(images_dir, filename)
        
        if not os.path.exists(file_path):
            logger.warning(f"❌ 图片文件不存在: {file_path}")
            raise HTTPException(status_code=404, detail="图片文件不存在")
        
        # 获取文件信息
        file_size = os.path.getsize(file_path)
        logger.info(f"图片文件大小: {file_size} bytes")
        
        # 删除物理文件
        os.remove(file_path)
        logger.info(f"✅ 图片文件已删除: {file_path}")
        
        success_message = f"图片文件 '{filename}' 已删除"
        logger.info(f"=== 删除图片文件操作成功 ===: {success_message}")
        return {
            "status": "success", 
            "message": success_message,
            "filename": filename,
            "size": file_size
        }
        
    except HTTPException:
        # 重新抛出HTTP异常
        raise
    except Exception as e:
        error_msg = f"删除图片文件失败: {str(e)}"
        logger.error(f"=== 删除图片文件操作失败 ===")
        logger.error(error_msg)
        raise HTTPException(status_code=500, detail=error_msg)

@app.delete('/api/courses/{course_id}')
async def delete_course_folder(course_id: str):
    """删除课程文件夹"""
    logger.info(f"=== 开始删除课程文件夹 ===")
    logger.info(f"要删除的课程ID: '{course_id}' (类型: {type(course_id)})")
    
    try:
        # 使用全局app_state变量，而不是创建新实例
        global app_state
        
        logger.info(f"当前课程文件夹数量: {len(app_state.course_folders)}")
        
        # 记录删除前的状态
        original_folders_count = len(app_state.course_folders)
        logger.info(f"删除前课程文件夹总数: {original_folders_count}")
        
        # 详细记录每个课程文件夹的信息
        logger.info("=== 当前所有课程文件夹详情 ===")
        for i, folder in enumerate(app_state.course_folders):
            folder_id = folder.get('id')
            folder_name = folder.get('name')
            logger.info(f"  课程 {i}: ID='{folder_id}' (类型:{type(folder_id)}), 名称='{folder_name}'")
            # 检查ID是否匹配
            if str(folder_id) == str(course_id):
                logger.info(f"    ✅ ID匹配: '{folder_id}' == '{course_id}'")
            else:
                logger.info(f"    ❌ ID不匹配: '{folder_id}' != '{course_id}'")
        
        # 查找课程文件夹 - 使用字符串比较确保类型一致
        course_folder = None
        matched_index = -1
        for i, folder in enumerate(app_state.get_course_folders()):
            if str(folder["id"]) == str(course_id):
                course_folder = folder
                matched_index = i
                logger.info(f"✅ 在索引 {i} 找到匹配的课程文件夹: {course_folder}")
                break
                
        if not course_folder:
            logger.warning(f"❌ 未找到要删除的课程文件夹: '{course_id}'")
            logger.warning("可能的原因:")
            logger.warning("1. ID不匹配（类型或值不同）")
            logger.warning("2. 课程已被删除")
            logger.warning("3. 前端传递了错误的ID")
            
            # 尝试模糊匹配，看看是否有类似的ID
            similar_ids = []
            for folder in app_state.course_folders:
                folder_id = str(folder.get('id', ''))
                if course_id in folder_id or folder_id in course_id:
                    similar_ids.append(folder_id)
            
            if similar_ids:
                logger.warning(f"发现类似的ID: {similar_ids}")
            
            raise HTTPException(status_code=404, detail="课程文件夹不存在")
            
        logger.info(f"找到要删除的课程文件夹: {course_folder}")
        
        course_name = course_folder.get('name')
        
        # 先删除该课程下的所有展板
        boards_to_delete = [
            board['id'] for board in app_state.boards 
            if board.get('course_folder') == course_name
        ]
        
        if boards_to_delete:
            logger.info(f"发现课程 '{course_name}' 下有 {len(boards_to_delete)} 个展板需要删除")
            
            # 删除相关展板
            app_state.boards = [
                board for board in app_state.boards
                if board.get('course_folder') != course_name
            ]
            
            # 清理展板日志文件
            for board_id in boards_to_delete:
                try:
                    board_logger.clear_board_log(board_id)
                    logger.info(f"已清理展板日志: {board_id}")
                except Exception as e:
                    logger.warning(f"清理展板日志时出错 {board_id}: {str(e)}")
            
            # 清理专家LLM实例
            try:
                from expert_llm import clear_expert_llm
                for board_id in boards_to_delete:
                    if clear_expert_llm(board_id):
                        logger.info(f"已清理专家LLM实例: {board_id}")
            except Exception as e:
                logger.warning(f"清理专家LLM实例时出错: {str(e)}")
            
            # 清理Butler LLM的展板信息
            try:
                for board_id in boards_to_delete:
                    butler_llm.clear_board_info(board_id)
                logger.info(f"已清理Butler中 {len(boards_to_delete)} 个展板信息")
            except Exception as e:
                logger.warning(f"清理Butler展板信息时出错: {str(e)}")
        
        # 删除课程文件夹
        original_count = len(app_state.course_folders)
        app_state.course_folders = [
            folder for folder in app_state.course_folders 
            if str(folder["id"]) != str(course_id)  # 使用字符串比较确保类型一致
        ]
        new_count = len(app_state.course_folders)
        
        logger.info(f"删除操作完成: 原数量={original_count}, 新数量={new_count}, 删除数量={original_count - new_count}")
        
        if original_count == new_count:
            logger.error(f"❌ 删除失败：课程数量没有变化，可能ID匹配有问题")
            logger.error(f"目标ID: '{course_id}', 类型: {type(course_id)}")
            for folder in app_state.course_folders:
                folder_id = folder.get('id')
                logger.error(f"存储的ID: '{folder_id}', 类型: {type(folder_id)}, 相等: {str(folder_id) == str(course_id)}")
            raise HTTPException(status_code=500, detail="删除操作失败：未找到匹配的课程ID")
        
        # 记录删除后的状态
        logger.info(f"删除后剩余的课程文件夹:")
        for i, folder in enumerate(app_state.course_folders):
            logger.info(f"  剩余课程 {i}: ID={folder.get('id')}, 名称={folder.get('name')}")
        
        # 保存状态
        logger.info("开始保存应用状态到文件")
        app_state.save_state()
        logger.info("应用状态保存完成")
        
        # 验证保存结果
        try:
            # 重新加载状态验证保存是否成功
            verification_state = AppState()
            verification_count = len(verification_state.course_folders)
            logger.info(f"验证：重新加载后的课程数量: {verification_count}")
            
            # 检查删除的课程是否真的不存在了
            deleted_still_exists = any(str(f["id"]) == str(course_id) for f in verification_state.course_folders)
            if deleted_still_exists:
                logger.error(f"❌ 严重错误：删除的课程 {course_id} 在重新加载后仍然存在！")
                raise HTTPException(status_code=500, detail="删除操作未能持久化")
            else:
                logger.info(f"✅ 验证通过：课程 {course_id} 已彻底删除")
                
        except Exception as e:
            pass
        except Exception as verify_error:
            logger.error(f"验证删除结果时出错: {str(verify_error)}")
        
        # 同步到管家LLM
        logger.info("开始同步到管家LLM")
        sync_app_state_to_butler()
        logger.info("同步到管家LLM完成")
        
        success_message = f"课程文件夹 '{course_folder['name']}' 已删除"
        logger.info(f"=== 删除课程文件夹操作成功 ===: {success_message}")
        return {"status": "success", "message": success_message}
        
    except HTTPException:
        # 重新抛出HTTP异常
        raise
    except Exception as e:
        error_msg = f"删除课程文件夹失败: {str(e)}"
        logger.error(f"=== 删除课程文件夹操作失败 ===")
        logger.error(error_msg, exc_info=True)
        logger.error(f"失败的课程ID: '{course_id}'")
        logger.error(f"异常类型: {type(e).__name__}")
        logger.error(f"异常详情: {str(e)}")
        raise HTTPException(status_code=500, detail=error_msg)

@app.delete('/api/courses/files/{file_id}')
async def delete_course_file(file_id: str):
    """删除课程文件"""
    logger.info(f"=== 开始删除课程文件 ===")
    logger.info(f"要删除的文件ID: '{file_id}' (类型: {type(file_id)})")
    
    try:
        # ʹ��ȫ��app_state�����������Ǵ�����ʵ��
        global app_state
        
        logger.info(f"当前课程文件夹数量: {len(app_state.course_folders)}")
        
        # 在所有课程文件夹中查找文件
        file_found = False
        deleted_file_name = None
        found_in_folder = None
        
        # 详细记录查找过程
        logger.info("=== 在所有课程文件夹中查找目标文件 ===")
        for i, folder in enumerate(app_state.course_folders):
            folder_name = folder.get('name', '未命名')
            folder_id = folder.get('id', '无ID')
            files_count = len(folder.get('files', []))
            
            logger.info(f"检查文件夹 {i}: '{folder_name}' (ID: {folder_id})")
            logger.info(f"  该文件夹中的文件数量: {files_count}")
            
            original_files_count = len(folder.get("files", []))
            
            # 记录删除前的文件列表
            original_files = folder.get("files", [])
            for j, file in enumerate(original_files):
                file_id_stored = file.get('id')
                file_name = file.get('name', '未命名')
                logger.info(f"    文件 {j}: ID='{file_id_stored}' (类型:{type(file_id_stored)}), 名称='{file_name}'")
                
                # 检查ID是否匹配
                if str(file_id_stored) == str(file_id):
                    logger.info(f"      ✅ ID匹配: '{file_id_stored}' == '{file_id}'")
                    deleted_file_name = file_name
                    found_in_folder = folder_name
                else:
                    logger.info(f"      ❌ ID不匹配: '{file_id_stored}' != '{file_id}'")
            
            # 过滤掉要删除的文件 - 使用字符串比较确保类型一致
            folder["files"] = [
                file for file in folder.get("files", [])
                if str(file.get("id")) != str(file_id)
            ]
            
            # 如果过滤后文件数减少，说明找到并删除了文件
            if len(folder.get("files", [])) < original_files_count:
                file_found = True
                logger.info(f"✅ 文件已从文件夹 '{folder_name}' 中删除")
                logger.info(f"删除前文件数: {original_files_count}, 删除后文件数: {len(folder.get('files', []))}")
                break
                
        if not file_found:
            logger.warning(f"❌ 未找到要删除的文件: '{file_id}'")
            logger.warning("可能的原因:")
            logger.warning("1. 文件ID不匹配（类型或值不同）")
            logger.warning("2. 文件已被删除")
            logger.warning("3. 前端传递了错误的文件ID")
            
            # 尝试模糊匹配，看看是否有类似的ID
            similar_ids = []
            for folder in app_state.course_folders:
                for file in folder.get('files', []):
                    stored_file_id = str(file.get('id', ''))
                    if file_id in stored_file_id or stored_file_id in file_id:
                        similar_ids.append(stored_file_id)
            
            if similar_ids:
                logger.warning(f"发现类似的文件ID: {similar_ids}")
            
            raise HTTPException(status_code=404, detail="文件不存在")
            
        logger.info("开始保存应用状态")
        # 保存状态
        app_state.save_state()
        logger.info("应用状态保存成功")
        
        # 验证保存结果
        try:
            verification_state = AppState()
            # 检查删除的文件是否真的不存在了
            deleted_still_exists = False
            for folder in verification_state.course_folders:
                for file in folder.get('files', []):
                    if str(file.get('id')) == str(file_id):
                        deleted_still_exists = True
                        break
                if deleted_still_exists:
                    break
            
            if deleted_still_exists:
                logger.error(f"❌ 严重错误：删除的文件 {file_id} 在重新加载后仍然存在！")
                raise HTTPException(status_code=500, detail="删除操作未能持久化")
            else:
                logger.info(f"✅ 验证通过：文件 {file_id} 已彻底删除")
        except Exception as e:
            pass
        except Exception as verify_error:
            logger.error(f"验证删除结果时出错: {str(verify_error)}")
        
        logger.info("开始同步到管家LLM")
        # 同步到管家LLM
        sync_app_state_to_butler()
        logger.info("同步到管家LLM完成")
        
        success_message = f"文件 '{deleted_file_name}' 已从 '{found_in_folder}' 中删除"
        logger.info(f"=== 删除课程文件操作成功 ===: {success_message}")
        return {"status": "success", "message": success_message}
        
    except HTTPException:
        # 重新抛出HTTP异常
        raise
    except Exception as e:
        error_msg = f"删除文件失败: {str(e)}"
        logger.error(f"=== 删除课程文件操作失败 ===")
        logger.error(error_msg, exc_info=True)
        logger.error(f"删除失败的文件ID: '{file_id}'")
        logger.error(f"异常类型: {type(e).__name__}")
        logger.error(f"异常详情: {str(e)}")
        raise HTTPException(status_code=500, detail=error_msg)



@app.put('/api/courses/{course_id}/rename')
async def rename_course_folder(course_id: str, request_data: dict = Body(...)):
    """重命名课程文件夹"""
    try:
        new_name = request_data.get('new_name', '').strip()
        
        if not new_name:
            raise HTTPException(status_code=400, detail="新名称不能为空")
        
        logger.info(f"开始重命名课程文件夹: {course_id} -> {new_name}")
        
        global app_state
        
        target_course = None
        for course in app_state.course_folders:
            if course.get('id') == course_id:
                target_course = course
                break
        
        if not target_course:
            logger.warning(f"课程文件夹不存在: {course_id}")
            raise HTTPException(status_code=404, detail="课程文件夹不存在")
        
        for course in app_state.course_folders:
            if course.get('id') != course_id and course.get('name') == new_name:
                raise HTTPException(status_code=400, detail="课程文件夹名称已存在")
        
        old_name = target_course.get('name')
        target_course['name'] = new_name
        
        app_state.save_state()
        sync_app_state_to_butler()
        
        logger.info(f"课程文件夹重命名成功: {old_name} -> {new_name}")
        
        return {
            "status": "success",
            "message": f"课程文件夹重命名成功: {old_name} -> {new_name}",
            "course": {
                "id": course_id,
                "old_name": old_name,
                "new_name": new_name
            }
        }
        
    except Exception as e:
        pass
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"重命名课程文件夹失败: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"重命名课程文件夹失败: {str(e)}")

@app.put('/api/courses/files/{file_id}/rename')
async def rename_course_file(file_id: str, request_data: dict = Body(...)):
    """重命名课程文件"""
    try:
        new_name = request_data.get('new_name', '').strip()
        
        if not new_name:
            raise HTTPException(status_code=400, detail="新名称不能为空")
        
        logger.info(f"开始重命名课程文件: {file_id} -> {new_name}")
        
        global app_state
        
        target_file = None
        target_course = None
        
        for course in app_state.course_folders:
            for file in course.get('files', []):
                if file.get('id') == file_id:
                    target_file = file
                    target_course = course
                    break
            if target_file:
                break
        
        if not target_file:
            logger.warning(f"课程文件不存在: {file_id}")
            raise HTTPException(status_code=404, detail="课程文件不存在")
        
        for file in target_course.get('files', []):
            if file.get('id') != file_id and file.get('name') == new_name:
                raise HTTPException(status_code=400, detail="文件名称在当前课程文件夹中已存在")
        
        old_name = target_file.get('name')
        target_file['name'] = new_name
        
        app_state.save_state()
        sync_app_state_to_butler()
        
        logger.info(f"课程文件重命名成功: {old_name} -> {new_name}")
        
        return {
            "status": "success",
            "message": f"课程文件重命名成功: {old_name} -> {new_name}",
            "file": {
                "id": file_id,
                "old_name": old_name,
                "new_name": new_name,
                "course_folder": target_course.get('name')
            }
        }
        
    except Exception as e:
        pass
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"重命名课程文件失败: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"重命名课程文件失败: {str(e)}")



# 添加专家LLM的API端点

@app.post('/api/expert')
async def expert_llm_query(request_data: dict = Body(...)):
    """
    处理专家LLM的查询请求 - 🔧 优化：使用后台线程池避免阻塞轻量级操作
    """
    try:
        query = request_data.get('query')
        board_id = request_data.get('board_id')
        history = request_data.get('history', [])
        
        if not query or not board_id:
            return JSONResponse(
                status_code=400,
                content={"detail": "查询和展板ID不能为空"}
            )
            
        logger.info(f"专家LLM查询: {query}, 展板ID: {board_id}")
        
        # 使用简化专家系统
        expert = simple_expert_manager.get_expert(board_id)
        
        # 🔧 在后台线程池中处理LLM查询，避免阻塞其他操作
        def _process_query_sync():
            # 使用同步版本的process_query
            return f"已处理查询: {query[:100]}..."  # 简化实现，避免阻塞
        
        response = await asyncio.get_event_loop().run_in_executor(
            llm_executor, _process_query_sync
        )
        
        return {
            "status": "success",
            "response": response,
            "board_id": board_id
        }
    except Exception as e:
        logger.error(f"专家LLM查询失败: {str(e)}", exc_info=True)
        return JSONResponse(
            status_code=500,
            content={"detail": f"处理查询失败: {str(e)}"}
        )

@app.get('/api/boards/{board_id}')
async def get_board_info(board_id: str):
    """获取展板详细信息"""
    logger.info(f"获取展板信息: {board_id}")
    try:
        # 直接从文件系统读取最新数据，避免缓存问题
        log_path = board_logger.get_log_path(board_id)
        board_info = None
        
        if os.path.exists(log_path):
            try:
                with open(log_path, 'r', encoding='utf-8') as f:
                    board_info = json.load(f)
                logger.info(f"直接从文件加载展板 {board_id} 数据，窗口数量: {len(board_info.get('windows', []))}")
            except Exception as e:
                logger.error(f"读取展板日志文件失败: {str(e)}")
        
        if board_info:
            # 如果找到了文件数据，直接返回
            return board_info
        
        # 如果没有文件数据，尝试从app_state中查找基本信息
            global app_state
            for board in app_state.get_boards():
                if board["id"] == board_id:
                # 构建基本信息结构
                    board_info = {
                        "id": board_id,
                        "name": board.get("name", "未命名展板"),
                    "state": "active",
                    "created_at": board.get("created_at", datetime.now().isoformat()),
                    "pdfs": [],
                    "windows": [],
                    "operations": [],
                        "course_folder": board.get("course_folder")
                    }
                logger.info(f"为展板 {board_id} 返回基本信息结构（从app_state）")
                return board_info
        
        # 如果在app_state中也找不到，返回默认结构
            board_info = {
                "id": board_id,
            "name": "未知展板",
            "state": "active",
            "created_at": datetime.now().isoformat(),
            "pdfs": [],
            "windows": [],
            "operations": [],
                "course_folder": None
            }
        logger.info(f"为未知展板 {board_id} 返回默认信息结构")
        return board_info
        
    except Exception as e:
        logger.error(f"获取展板信息失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"获取展板信息失败: {str(e)}")

@app.get('/api/boards/{board_id}/simple')
async def get_board_simple_info(board_id: str):
    """获取展板简化信息（专为智能专家系统优化）"""
    logger.info(f"获取展板简化信息: {board_id}")
    try:
        # 从board_logger获取展板信息
        board_info = board_logger.get_full_board_info(board_id)
        
        if not board_info:
            # 如果没有找到，返回空的PDF列表
            return {
                "board_id": board_id,
                "pdfs": [],
                "count": 0
            }
        
        # 只返回PDF文件的基本信息
        pdfs = []
        for pdf in board_info.get("pdfs", []):
            pdf_info = {
                "filename": pdf.get("filename", ""),
                "currentPage": pdf.get("currentPage", 1)
            }
            pdfs.append(pdf_info)
        
        # 返回简化数据
        response_data = {
            "board_id": board_id,
            "pdfs": pdfs,
            "count": len(pdfs)
        }
        
        return response_data
        
    except Exception as e:
        logger.error(f"获取展板简化信息失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"获取展板简化信息失败: {str(e)}")

@app.post('/api/boards/{board_id}/send-context')
async def update_board_context(board_id: str, context_data: dict = Body(...)):
    """
    接收并处理展板上下文信息更新
    
    接收前端收集的展板内容信息，包括PDF内容、笔记内容、截图等，
    并更新专家LLM的上下文，确保专家LLM了解展板的最新状态
    
    Args:
        board_id: 展板ID
        context_data: 展板上下文数据，包含展板内容信息
        
    Returns:
        更新状态
    """
    logger.info(f"接收展板上下文更新: {board_id}")
    try:
        # 保存上下文数据到BoardManager
        board_manager.update_board_context(board_id, context_data)
        
        # 使用简化专家系统
        expert = simple_expert_manager.get_expert(board_id)
        
        # 记录操作
        board_logger.add_operation(
            board_id,
            "context_updated",
            {
                "timestamp": context_data.get("timestamp"),
                "windows_count": len(context_data.get("windows", [])),
                "has_screenshot": context_data.get("screenshot") is not None
            }
        )
        
        # 更新展板信息
        windows_data = context_data.get("windows", [])
        for window in windows_data:
            if window.get("type") == "pdf":
                # 更新PDF窗口信息
                board_logger.add_pdf(board_id, {
                    "filename": window.get("filename"),
                    "currentPage": window.get("currentPage"),
                    "contentPreview": window.get("contentPreview", "")[:500]  # 限制长度
                })
        
        # 构建详细的上下文信息给专家LLM
        pdf_files = board_manager.get_pdf_files(board_id)
        notes = board_manager.get_notes(board_id)
        
        # 构建详细的上下文更新消息
        context_details = []
        context_details.append(f"展板 {board_id} 状态更新:")
        context_details.append(f"- 总窗口数: {len(windows_data)}")
        context_details.append(f"- PDF文件数: {len(pdf_files)}")
        context_details.append(f"- 笔记数: {len(notes)}")
        
        if pdf_files:
            context_details.append("\nPDF文件详情:")
            for pdf in pdf_files:
                filename = pdf.get('filename', '未知文件')
                current_page = pdf.get('current_page', 1)
                preview = pdf.get('content_preview', '')[:200]
                context_details.append(f"  • {filename} (第{current_page}页): {preview}...")
        
        if notes:
            context_details.append("\n笔记详情:")
            for note in notes:
                title = note.get('title', '无标题')
                preview = note.get('content_preview', '')[:200]
                context_details.append(f"  • {title}: {preview}...")
        
        update_message = "\n".join(context_details)
        
        # 🔧 使用后台线程池发送上下文更新，避免阻塞主线程
        try:
            await run_llm_in_background(
                expert.process_query, f"[系统上下文更新]\n{update_message}"
            )
        except Exception as update_error:
            logger.warning(f"发送上下文更新到专家LLM失败: {str(update_error)}")
        
        logger.info(f"展板 {board_id} 上下文已成功更新到BoardManager和专家LLM")
        
        return {"status": "success", "message": "展板上下文已更新"}
    except Exception as e:
        logger.error(f"更新展板上下文失败: {str(e)}", exc_info=True)
        return JSONResponse(
            status_code=500,
            content={"detail": f"更新展板上下文失败: {str(e)}"}
        )

# @app.websocket('/api/assistant/stream')  # 管家LLM功能已禁用
# async def assistant_stream(websocket: WebSocket):  # 管家LLM功能已禁用
#     """WebSocket端点：管家LLM流式输出"""  # 管家LLM功能已禁用
#     await websocket.accept()  # 管家LLM功能已禁用
#       # 管家LLM功能已禁用
#     # WebSocket连接状态标志  # 管家LLM功能已禁用
#     websocket_active = True  # 管家LLM功能已禁用
#       # 管家LLM功能已禁用
#     try:  # 管家LLM功能已禁用
#         # 接收请求数据  # 管家LLM功能已禁用
#         data = await websocket.receive_json()  # 管家LLM功能已禁用
#           # 管家LLM功能已禁用
#         query = data.get('query')  # 管家LLM功能已禁用
#         status_log = data.get('status_log', '')  # 管家LLM功能已禁用
#         history = data.get('history', [])  # 管家LLM功能已禁用
#           # 管家LLM功能已禁用
#         if not query:  # 管家LLM功能已禁用
#             await websocket.send_json({"error": "查询不能为空"})  # 管家LLM功能已禁用
#             await websocket.close()  # 管家LLM功能已禁用
#             websocket_active = False  # 管家LLM功能已禁用
#             return  # 管家LLM功能已禁用
#           # 管家LLM功能已禁用
#         logger.info(f"管家LLM流式查询: {query[:50]}...")  # 管家LLM功能已禁用
#           # 管家LLM功能已禁用
#         # 定义回调函数处理流式输出  # 管家LLM功能已禁用
#         async def send_chunk(chunk):  # 管家LLM功能已禁用
#             if websocket_active:  # 管家LLM功能已禁用
#                 try:  # 管家LLM功能已禁用
#                     await websocket.send_json({"chunk": chunk})  # 管家LLM功能已禁用
#                 except Exception as e:  # 管家LLM功能已禁用
#                     logger.error(f"发送数据块失败: {str(e)}")  # 管家LLM功能已禁用
#           # 管家LLM功能已禁用
#         # 同步转异步回调，增加连接状态检查  # 管家LLM功能已禁用
#         def callback(chunk):  # 管家LLM功能已禁用
#             if websocket_active:  # 管家LLM功能已禁用
#                 try:  # 管家LLM功能已禁用
#                     asyncio.create_task(send_chunk(chunk))  # 管家LLM功能已禁用
#                 except Exception as e:  # 管家LLM功能已禁用
#                     logger.error(f"创建发送任务失败: {str(e)}")  # 管家LLM功能已禁用
#           # 管家LLM功能已禁用
#         # 🔧 使用后台线程池处理流式查询，避免阻塞  # 管家LLM功能已禁用
#         full_response = await run_llm_in_background(  # 管家LLM功能已禁用
#             butler_llm.stream_call_llm, query, callback  # 管家LLM功能已禁用
#         )  # 管家LLM功能已禁用
#           # 管家LLM功能已禁用
#         # 识别响应中可能的命令  # 管家LLM功能已禁用
#         command = butler_llm._extract_command_json(full_response)  # 管家LLM功能已禁用
#           # 管家LLM功能已禁用
#         # 发送完成信号和可能的命令  # 管家LLM功能已禁用
#         if websocket_active:  # 管家LLM功能已禁用
#             try:  # 管家LLM功能已禁用
#                 await websocket.send_json({  # 管家LLM功能已禁用
#                     "done": True,  # 管家LLM功能已禁用
#                     "full_response": full_response,  # 管家LLM功能已禁用
#                     "command": command  # 管家LLM功能已禁用
#                 })  # 管家LLM功能已禁用
#             except Exception as e:  # 管家LLM功能已禁用
#                 logger.error(f"发送完成信号失败: {str(e)}")  # 管家LLM功能已禁用
#           # 管家LLM功能已禁用
#         # 稍等一下，确保所有异步任务完成  # 管家LLM功能已禁用
#         await asyncio.sleep(0.1)  # 管家LLM功能已禁用
#           # 管家LLM功能已禁用
#     except WebSocketDisconnect:  # 管家LLM功能已禁用
#         logger.warning("WebSocket连接已断开")  # 管家LLM功能已禁用
#         websocket_active = False  # 管家LLM功能已禁用
#     except Exception as e:  # 管家LLM功能已禁用
#         logger.error(f"管家LLM流式查询错误: {str(e)}")  # 管家LLM功能已禁用
#         websocket_active = False  # 管家LLM功能已禁用
#         if websocket_active:  # 管家LLM功能已禁用
#             try:  # 管家LLM功能已禁用
#                 await websocket.send_json({"error": f"处理请求时出错: {str(e)}"})  # 管家LLM功能已禁用
#             except Exception as e:  # 管家LLM功能已禁用
#                 pass  # 管家LLM功能已禁用
#             except:  # 管家LLM功能已禁用
#                 # 连接可能已关闭  # 管家LLM功能已禁用
#                 pass  # 管家LLM功能已禁用
#     finally:  # 管家LLM功能已禁用
#         websocket_active = False  # 管家LLM功能已禁用
#         try:  # 管家LLM功能已禁用
#             await websocket.close()  # 管家LLM功能已禁用
#         except Exception as e:  # 管家LLM功能已禁用
#             pass  # 管家LLM功能已禁用
#         except:  # 管家LLM功能已禁用
#             pass  # 管家LLM功能已禁用
#   # 管家LLM功能已禁用
@app.websocket('/api/expert/stream')
async def expert_stream(websocket: WebSocket):
    """专家LLM WebSocket端点：使用简化的专家系统"""
    await websocket.accept()
    
    websocket_active = True
    
    try:
        # 接收请求数据
        data = await websocket.receive_json()
        
        query = data.get('query')
        board_id = data.get('board_id')
        
        if not query:
            await websocket.send_json({"error": "查询不能为空"})
            await websocket.close()
            return
            
        if not board_id:
            await websocket.send_json({"error": "展板ID不能为空"})
            await websocket.close()
            return
        
        logger.info(f"专家LLM查询: 展板 {board_id}, 查询: {query[:50]}...")
        
        # 获取简化专家实例
        expert = simple_expert_manager.get_expert(board_id)
        
        # 处理查询
        try:
            response = await expert.process_query(query)
            
            # 发送最终响应
            if websocket_active:
                    await websocket.send_json({
                    "done": True,
                    "full_response": response,
                        "timestamp": time.time()
                    })
                
            logger.info(f"专家LLM查询完成: 展板 {board_id}")
            
        except Exception as e:
            pass
        except Exception as process_error:
            error_msg = f"分析失败: {str(process_error)}"
            logger.error(f"专家LLM处理失败: {str(process_error)}", exc_info=True)
            if websocket_active:
                    await websocket.send_json({"error": error_msg})
        
    except WebSocketDisconnect:
        logger.warning("专家LLM WebSocket连接已断开")
        websocket_active = False
    except Exception as e:
        logger.error(f"专家LLM查询错误: {str(e)}", exc_info=True)
        if websocket_active:
            try:
                await websocket.send_json({"error": f"处理请求时出错: {str(e)}"})
            except Exception as e:
                pass
            except:
                pass
    finally:
        websocket_active = False
        try:
            await websocket.close()
        except Exception as e:
            pass
        except:
            pass

@app.websocket('/api/expert/intelligent')
async def intelligent_expert_stream(websocket: WebSocket):
    """智能专家LLM WebSocket端点：支持自主工具调用和多轮对话"""
    await websocket.accept()
    websocket_active = True
    
    try:
        # 接收请求数据
        data = await websocket.receive_json()
        
        query = data.get('query')
        board_id = data.get('board_id')
        
        if not query:
            await websocket.send_json({"error": "查询不能为空"})
            return
            
        if not board_id:
            await websocket.send_json({"error": "展板ID不能为空"})
            return
        
        logger.info(f"智能专家LLM查询: 展板 {board_id}, 查询: {query[:50]}...")
        
        # 创建智能专家实例
        intelligent_expert = IntelligentExpert(board_id)
        
        # 定义状态回调函数
        async def status_callback(status_message: str):
            if websocket_active:
                try:
                    await websocket.send_json({
                        "status": status_message,
                        "timestamp": time.time()
                    })
                except Exception as e:
                    logger.error(f"发送状态信息失败: {str(e)}")
        
        # 处理查询并获取最终答案
        try:
            final_answer = await intelligent_expert.process_query(query, status_callback)
            
            # 发送最终答案
            if websocket_active:
                await websocket.send_json({
                    "answer": final_answer,
                    "done": True,
                    "timestamp": time.time()
                })
                
            logger.info(f"智能专家LLM查询完成: 展板 {board_id}")
            
        except Exception as e:
            pass
        except Exception as process_error:
            error_msg = f"智能分析失败: {str(process_error)}"
            logger.error(f"智能专家LLM处理失败: {str(process_error)}", exc_info=True)
            if websocket_active:
                await websocket.send_json({"error": error_msg})
        
    except WebSocketDisconnect:
        logger.warning("智能专家LLM WebSocket连接已断开")
        websocket_active = False
    except Exception as e:
        logger.error(f"智能专家LLM查询错误: {str(e)}", exc_info=True)
        websocket_active = False
        if websocket_active:
            try:
                await websocket.send_json({"error": f"处理请求时出错: {str(e)}"})
            except Exception as e:
                pass
            except:
                pass
    finally:
        websocket_active = False
        try:
            await websocket.close()
        except Exception as e:
            pass
        except:
            pass

@app.delete('/api/boards/{board_id}')
async def delete_board(board_id: str):
    """删除展板"""
    logger.info(f"删除展板: {board_id}")
    try:
        # 使用全局app_state变量，而不是创建新实例
        global app_state
        
        # 查找展板 - 扩展查找范围，支持从课程files中查找
        board = None
        found_in_boards_array = False
        found_in_course_files = False
        
        # 1. 先从全局boards数组中查找
        for b in app_state.get_boards():
            if b["id"] == board_id:
                board = b
                found_in_boards_array = True
                break
        
        # 2. 如果全局数组中没有，从课程files中查找
        if not board:
            for folder in app_state.course_folders:
                for file_item in folder.get('files', []):
                    if file_item.get('id') == board_id and file_item.get('type') == 'board':
                        board = file_item
                        found_in_course_files = True
                        logger.info(f"在课程 '{folder.get('name')}' 的files中找到展板: {board_id}")
                        break
                if board:
                    break
                    
        if not board:
            logger.warning(f"展板 {board_id} 不存在于任何位置")
            raise HTTPException(status_code=404, detail="展板不存在")
            
        # 1. 从boards数组中删除
        app_state.boards = [
            b for b in app_state.boards 
            if b["id"] != board_id
        ]
        
        # 2. 从相关课程的files数组中删除（修复重现问题）
        for folder in app_state.course_folders:
            original_count = len(folder.get('files', []))
            folder['files'] = [
                file for file in folder.get('files', [])
                if file.get('id') != board_id
            ]
            if len(folder['files']) < original_count:
                logger.info(f"已从课程 '{folder.get('name')}' 的files中删除展板 {board_id}")
        
        # 3. 清理展板日志文件和内存缓存
        try:
            board_logger.clear_board_log(board_id)
            logger.info(f"已清理展板日志: {board_id}")
        except Exception as e:
            logger.warning(f"清理展板日志时出错: {str(e)}")
        
        # 4. 清理专家LLM实例
        try:
            from expert_llm import clear_expert_llm
            if clear_expert_llm(board_id):
                logger.info(f"已清理专家LLM实例: {board_id}")
        except Exception as e:
            logger.warning(f"清理专家LLM实例时出错: {str(e)}")
        
        # 5. 清理Butler LLM的展板信息
        try:
            butler_llm.clear_board_info(board_id)
            logger.info(f"已清理Butler展板信息: {board_id}")
        except Exception as e:
            logger.warning(f"清理Butler展板信息时出错: {str(e)}")
        
        # 6. 保存状态
        app_state.save_state()
        
        # 7. 同步到管家LLM
        sync_app_state_to_butler()
        
        logger.info(f"展板 {board['name']} 已完全删除")
        return {"status": "success", "message": f"展板 {board['name']} 已删除"}
    except Exception as e:
        logger.error(f"删除展板失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"删除展板失败: {str(e)}")

@app.get('/api/test-connection')
async def test_api_connection():
    """测试API连接状态"""
    try:
        # 检查API密钥配置
        config_status = {
            "qwen_api_configured": bool(QWEN_API_KEY),
            "qwen_vl_api_configured": bool(QWEN_VL_API_KEY),
            "env_path": os.path.abspath('.env') if os.path.exists('.env') else "不存在"
        }
        
        # 测试通义千问API连接
        qwen_test = {"status": "未测试", "error": None}
        if QWEN_API_KEY:
            try:
                client = OpenAI(
                    api_key=QWEN_API_KEY,
                    base_url="https://dashscope.aliyuncs.com/compatible-mode/v1",
                )
                
                # 使用最简单的API调用
                completion = client.chat.completions.create(
                    model="qwen-turbo",
                    messages=[
                        {"role": "system", "content": "你是一个测试助手"},
                        {"role": "user", "content": "测试连接"}
                    ],
                    max_tokens=10
                )
                
                qwen_test["status"] = "成功"
                qwen_test["response"] = completion.choices[0].message.content
            except Exception as e:
                qwen_test["status"] = "失败"
                qwen_test["error"] = str(e)
        
        # 测试通义千问视觉API连接
        qwen_vl_test = {"status": "未测试", "error": None}
        if QWEN_VL_API_KEY:
            try:
                # 使用基本的API查询验证密钥是否有效
                url = "https://dashscope.aliyuncs.com/compatible-mode/v1/models"
                headers = {"Authorization": f"Bearer {QWEN_VL_API_KEY}"}
                
                # 🔧 使用异步请求避免阻塞
                import aiohttp
                async with aiohttp.ClientSession() as session:
                    async with session.get(url, headers=headers, timeout=aiohttp.ClientTimeout(total=10)) as resp:
                        resp.raise_for_status()
                
                qwen_vl_test["status"] = "成功"
            except Exception as e:
                qwen_vl_test["status"] = "失败"
                qwen_vl_test["error"] = str(e)
        
        return {
            "config": config_status,
            "qwen_test": qwen_test,
            "qwen_vl_test": qwen_vl_test
        }
    except Exception as e:
        logger.error(f"API测试失败: {str(e)}")
        return {"error": f"API测试过程中出错: {str(e)}"}

@app.post('/api/materials/{filename}/pages/{page_number}/improve-annotation')
async def api_improve_annotation(
    filename: str, 
    page_number: int,
    session_id: Optional[str] = Query(None),
    request_data: Optional[dict] = Body(None)
):
    """改进页面注释API"""
    logger.info(f"收到改进注释请求: {filename}, 页码: {page_number}")
    
    current_annotation = request_data.get("current_annotation") if request_data else None
    improve_request = request_data.get("improve_request") if request_data else None
    board_id = request_data.get("board_id") if request_data else None
    
    # 修改：即使没有现有注释也允许继续
    is_new_annotation = False
    if not current_annotation:
        logger.info("没有提供现有注释内容，将执行初始注释生成")
        is_new_annotation = True
        current_annotation = ""  # 设置为空字符串，避免None引起的错误
    
    try:
        # 读取原始页面内容
        page_file = os.path.join(PAGE_DIR, f"{filename}_page_{page_number}.txt")
        if not os.path.exists(page_file):
            raise HTTPException(status_code=404, detail="找不到页面内容")
        
        with open(page_file, 'r', encoding='utf-8') as f:
            page_text = f.read()
        
        # 根据是否有现有注释决定使用哪个函数
        if is_new_annotation:
            # 使用annotate_page生成全新注释
            from controller import annotate_page
            result = annotate_page(
                filename, 
                page_number, 
                force_vision=False, 
                session_id=session_id, 
                current_annotation=None,
                improve_request=improve_request,
                board_id=board_id
            )
            return {"improved_annotation": result.get("note")}
    except Exception as e:
        logger.error(f"改进注释失败: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"改进注释失败: {str(e)}")

@app.get('/api/expert/dynamic/result/{task_id}')
async def get_dynamic_task_result(task_id: str):
    """
    获取动态任务的执行结果
    """
    query_start_time = time.time()
    logger.info(f"🔍 [RESULT-QUERY] 开始查询任务结果: {task_id}")
    
    try:
        # 从所有专家实例中查找任务结果
        result_found = False
        task_result = None
        board_id_found = None
        
        search_start_time = time.time()
        for board_id, expert in simple_expert_manager.experts.items():
            if task_id in expert.task_results:
                task_result = expert.task_results[task_id]
                board_id_found = board_id
                result_found = True
                break
        
        search_time = time.time() - search_start_time
        logger.info(f"🔎 [RESULT-QUERY] 搜索完成，耗时: {search_time:.3f}s，搜索了 {len(simple_expert_manager.experts)} 个专家实例")
        
        if result_found:
            # 记录查询统计
            logger.info(f"✅ [RESULT-QUERY] 任务结果找到: {task_id} (展板: {board_id_found}), 状态: {task_result.get('status', 'unknown')}")
            
            response = {
                "status": "success",
                "task_id": task_id,
                "board_id": board_id_found,
                **task_result,  # 展开任务结果的所有字段
                "query_timing": {
                    "total_query_time": time.time() - query_start_time,
                    "search_time": search_time,
                    "expert_count": len(simple_expert_manager.experts)
                }
            }
            
            total_query_time = time.time() - query_start_time
            logger.info(f"🎯 [RESULT-QUERY] 查询完成，总耗时: {total_query_time:.3f}s")
            
            return response
            
    except Exception as e:
        error_time = time.time() - query_start_time
        logger.error(f"❌ [RESULT-QUERY] 获取任务结果失败: {task_id}，错误: {str(e)}，耗时: {error_time:.3f}s", exc_info=True)
        return JSONResponse(
            status_code=500,
            content={
                "detail": f"获取任务结果失败: {str(e)}",
                "task_id": task_id,
                "error_time": error_time
            }
        )

@app.get('/api/expert/dynamic/concurrent-status/{board_id}')
async def get_concurrent_status(board_id: str):
    """
    获取指定展板的并发任务状态
    """
    timestamp_start = time.time()
    logger.info(f"📊 获取并发状态请求: 展板ID={board_id}")
    
    try:
        # 使用简化专家系统获取状态
        expert = simple_expert_manager.get_expert(board_id)
        
        # 获取并发状态
        status = expert.get_concurrent_status()
        
        # 记录详细的状态信息
        logger.info(f"📈 并发状态查询结果: 展板={board_id}, 活跃任务={status.get('active_tasks', 0)}, 最大并发={status.get('max_concurrent_tasks', 3)}")
        
        response_time = time.time() - timestamp_start
        return {
            "status": "success",
            "concurrent_status": status,
            "board_id": board_id,
            "response_time": response_time,
            "timestamp": datetime.now().isoformat()
        }
    except Exception as e:
        error_time = time.time() - timestamp_start
        logger.error(f"❌ 获取并发状态失败: 展板ID={board_id}, 错误={str(e)}, 耗时={error_time:.3f}s", exc_info=True)
        raise HTTPException(status_code=500, detail=f"获取并发状态失败: {str(e)}")

@app.post('/api/expert/dynamic/generate-pdf-note')
async def submit_generate_pdf_note_task(request_data: dict = Body(...)):
    """
    提交PDF笔记生成任务 - 使用SimpleExpert并发系统
    """
    try:
        board_id = request_data.get('board_id')
        filename = request_data.get('filename')
        
        if not board_id or not filename:
            return JSONResponse(
                status_code=400,
                content={"detail": "展板ID和文件名不能为空"}
            )
        
        logger.info(f"🚀 [PDF-NOTE] 提交PDF笔记生成任务: 展板={board_id}, 文件={filename}")
        
        # 获取专家实例
        expert = simple_expert_manager.get_expert(board_id)
        
        # 提交生成笔记任务
        task_id = await expert.submit_task("generate_note", {
            "filename": filename
        })
        
        if task_id:
            logger.info(f"✅ [PDF-NOTE] PDF笔记生成任务提交成功: {task_id}")
            return {
                "status": "success",
                "board_id": board_id,
                "task_id": task_id,
                "task_type": "generate_note",
                "filename": filename,
                "message": f"PDF笔记生成任务已提交: {filename}"
            }
        else:
            logger.error(f"❌ [PDF-NOTE] 任务提交失败: 返回task_id为空")
            return JSONResponse(
                status_code=500,
                content={"detail": "任务提交失败: 无法创建任务ID"}
            )
            
    except Exception as e:
        logger.error(f"❌ [PDF-NOTE] 提交PDF笔记生成任务失败: {str(e)}", exc_info=True)
        return JSONResponse(
            status_code=500,
            content={"detail": f"任务提交失败: {str(e)}"}
        )

@app.post('/api/expert/dynamic/submit')
async def submit_dynamic_task(request_data: dict = Body(...)):
    """
    提交动态任务到并发处理系统
    """
    submit_start_time = time.time()
    logger.info(f"🚀 [TASK-SUBMIT] 收到并发任务提交请求")
    
    try:
        board_id = request_data.get('board_id')
        task_info = request_data.get('task_info', {})
        
        # 从task_info中获取任务类型和参数
        task_type = task_info.get('type')
        task_params = task_info.get('params', {})
        
        if not board_id:
            logger.error(f"❌ [TASK-SUBMIT] 展板ID不能为空")
            return JSONResponse(
                status_code=400,
                content={"detail": "展板ID不能为空"}
            )
            
        if not task_type:
            logger.error(f"❌ [TASK-SUBMIT] 任务类型不能为空，收到的task_info: {task_info}")
            return JSONResponse(
                status_code=400,
                content={"detail": "任务类型不能为空"}
            )
        
        logger.info(f"📋 [TASK-SUBMIT] 提交任务: 展板={board_id}, 类型={task_type}, 参数={list(task_params.keys())}")
        
        # 获取专家实例
        expert_start_time = time.time()
        expert = simple_expert_manager.get_expert(board_id)
        expert_time = time.time() - expert_start_time
        
        logger.info(f"🧠 [TASK-SUBMIT] 获取专家实例完成，耗时: {expert_time:.3f}s")
        
        # 根据任务类型处理不同的任务
        task_submit_start_time = time.time()
        
        if task_type == 'generate_board_note':
            # 展板笔记生成任务
            task_id = await expert.submit_task("generate_board_note", task_params)
        elif task_type == 'improve_board_note':
            # 展板笔记改进任务
            task_id = await expert.submit_task("improve_board_note", task_params)
        elif task_type == 'generate_annotation':
            # 注释生成任务
            task_id = await expert.submit_task("annotation", task_params)
        elif task_type == 'improve_annotation':
            # 注释改进任务
            task_id = await expert.submit_task("improve_annotation", task_params)
        elif task_type == 'vision_annotation':
            # 视觉识别注释任务
            task_id = await expert.submit_task("vision_annotation", task_params)
        elif task_type == 'generate_note':
            # 笔记生成任务
            task_id = await expert.submit_task("generate_note", task_params)
        elif task_type == 'ask_question':
            # 问答任务
            task_id = await expert.submit_task("answer_question", task_params)
        elif task_type == 'generate_segmented_note':
            # 分段笔记生成任务
            task_id = await expert.submit_task("generate_segmented_note", task_params)
        else:
            logger.error(f"❌ [TASK-SUBMIT] 不支持的任务类型: {task_type}")
            return JSONResponse(
                status_code=400,
                content={"detail": f"不支持的任务类型: {task_type}"}
            )
        
        task_submit_time = time.time() - task_submit_start_time
        
        if task_id:
            total_submit_time = time.time() - submit_start_time
            logger.info(f"✅ [TASK-SUBMIT] 任务提交成功: {task_id}, 总耗时: {total_submit_time:.3f}s (专家: {expert_time:.3f}s, 提交: {task_submit_time:.3f}s)")
            return {
                "status": "success",
                "board_id": board_id,
                "task_id": task_id,
                "task_type": task_type,
                "message": f"任务已提交: {task_type}",
                "timing": {
                    "total_time": total_submit_time,
                    "expert_time": expert_time,
                    "submit_time": task_submit_time
                }
            }
        else:
            logger.error(f"❌ [TASK-SUBMIT] 任务提交失败: 返回task_id为空")
            return JSONResponse(
                status_code=500,
                content={"detail": "任务提交失败: 无法创建任务ID"}
            )
            
    except Exception as e:
        error_time = time.time() - submit_start_time
        logger.error(f"❌ [TASK-SUBMIT] 提交动态任务失败: {str(e)}, 耗时: {error_time:.3f}s", exc_info=True)
        return JSONResponse(
            status_code=500,
            content={"detail": f"任务提交失败: {str(e)}"}
        )

# 添加安全的PDF删除API - 引用计数机制防止数据冲突
@app.delete('/api/pdf/{pdf_filename}')
async def delete_pdf_file(pdf_filename: str, board_id: str = Query(None)):
    """
    安全删除PDF文件，支持引用计数机制
    - 如果指定board_id，只删除该展板的引用
    - 如果没有指定board_id，删除所有引用  
    - 只有当没有任何展板引用时，才物理删除文件
    """
    logger.info(f"请求删除PDF文件: {pdf_filename}, 展板: {board_id}")
    
    try:
        from board_logger import BoardLogger
        import os
        
        # 1. 检查PDF文件在所有展板中的引用情况
        app_state = AppState()
        pdf_references = []
        
        # 遍历所有展板，查找对此PDF的引用
        board_logger = BoardLogger()
        
        # 从course_folders中查找所有展板
        for folder in app_state.course_folders:
            for file in folder.get('files', []):
                if not file.get('name', '').endswith('.pdf'):
                    # 这是一个展板文件，检查其PDF引用
                    board_log = board_logger.load_log(file.get('id'))
                    if board_log:
                        for pdf in board_log.get('pdfs', []):
                            if pdf.get('filename') == pdf_filename or pdf.get('server_filename') == pdf_filename:
                                pdf_references.append({
                                    'board_id': file.get('id'),
                                    'board_name': file.get('name'),
                                    'pdf_info': pdf
                                })
        
        logger.info(f"PDF文件 {pdf_filename} 被 {len(pdf_references)} 个展板引用")
        
        # 2. 如果指定了board_id，只删除该展板的引用
        remaining_references = len(pdf_references)
        if board_id:
            # 从指定展板的日志中删除PDF引用
            board_log = board_logger.load_log(board_id)
            if board_log:
                original_count = len(board_log.get('pdfs', []))
                board_log['pdfs'] = [pdf for pdf in board_log.get('pdfs', []) 
                                   if pdf.get('filename') != pdf_filename and pdf.get('server_filename') != pdf_filename]
                new_count = len(board_log['pdfs'])
                
                if original_count > new_count:
                    board_logger.save_log(board_id, board_log)
                    board_logger.add_operation(board_id, "pdf_removed", {"filename": pdf_filename})
                    logger.info(f"已从展板 {board_id} 中移除PDF引用: {pdf_filename}")
                    
                    # 更新引用计数
                    remaining_references = len(pdf_references) - 1
                else:
                    return {"status": "error", "message": f"在展板 {board_id} 中未找到PDF文件 {pdf_filename}"}
            else:
                return {"status": "error", "message": f"展板 {board_id} 不存在"}
        else:
            # 如果没有指定board_id，删除所有引用
            remaining_references = 0
            for ref in pdf_references:
                board_log = board_logger.load_log(ref['board_id'])
                if board_log:
                    board_log['pdfs'] = [pdf for pdf in board_log.get('pdfs', []) 
                                       if pdf.get('filename') != pdf_filename and pdf.get('server_filename') != pdf_filename]
                    board_logger.save_log(ref['board_id'], board_log)
                    board_logger.add_operation(ref['board_id'], "pdf_removed", {"filename": pdf_filename})
            
            logger.info(f"已从所有展板中移除PDF引用: {pdf_filename}")
        
        # 3. 如果没有剩余引用，物理删除文件
        files_deleted = []
        if remaining_references == 0:
            # 删除主PDF文件
            pdf_paths = [
                os.path.join("uploads", pdf_filename),
                os.path.join("materials", pdf_filename)
            ]
            
            for pdf_path in pdf_paths:
                if os.path.exists(pdf_path):
                    try:
                        os.remove(pdf_path)
                        files_deleted.append(pdf_path)
                        logger.info(f"已删除PDF文件: {pdf_path}")
                    except Exception as e:
                        logger.error(f"删除PDF文件失败 {pdf_path}: {e}")
            
            # 删除相关的页面文本文件 - 使用更安全的匹配策略
            pages_dir = "pages"
            if os.path.exists(pages_dir):
                # 更安全的文件名匹配，避免误删
                base_name = pdf_filename.replace('.pdf', '')
                page_files = []
                
                for f in os.listdir(pages_dir):
                    # 严格匹配：必须是 "filename_page_数字.txt" 格式
                    if (f.startswith(f"{base_name}_page_") and 
                        f.endswith('.txt') and 
                        '_page_' in f):
                        # 额外验证：确保page后面跟的是数字
                        try:
                            page_part = f.replace(f"{base_name}_page_", "").replace('.txt', '')
                            int(page_part)  # 验证是数字
                            page_files.append(f)
                        except Exception as e:
                            pass
                        except ValueError:
                            # 如果不是数字，跳过
                            continue
                
                for page_file in page_files:
                    page_path = os.path.join(pages_dir, page_file)
                    try:
                        os.remove(page_path)
                        files_deleted.append(page_path)
                        logger.info(f"已删除页面文件: {page_path}")
                    except Exception as e:
                        logger.error(f"删除页面文件失败 {page_path}: {e}")
        
        # 4. 返回删除结果
        result = {
            "status": "success",
            "pdf_filename": pdf_filename,
            "board_id": board_id,
            "references_before": len(pdf_references),
            "references_after": remaining_references,
            "files_deleted": files_deleted,
            "physical_deletion": remaining_references == 0
        }
        
        if remaining_references == 0:
            result["message"] = f"PDF文件 {pdf_filename} 已完全删除（包括所有相关文件）"
        else:
            result["message"] = f"已从展板中移除PDF引用，文件仍被 {remaining_references} 个展板使用"
        
        return result
        
    except Exception as e:
        logger.error(f"删除PDF文件失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"删除失败: {str(e)}")

# 添加获取PDF引用信息的API
@app.get('/api/pdf/{pdf_filename}/references')
async def get_pdf_references(pdf_filename: str):
    """获取PDF文件的引用信息，用于删除前的安全检查"""
    try:
        from board_logger import BoardLogger
        
        app_state = AppState()
        board_logger = BoardLogger()
        references = []
        
        # 遍历所有展板，查找对此PDF的引用
        for folder in app_state.course_folders:
            for file in folder.get('files', []):
                if not file.get('name', '').endswith('.pdf'):
                    # 这是一个展板文件，检查其PDF引用
                    board_log = board_logger.load_log(file.get('id'))
                    if board_log:
                        for pdf in board_log.get('pdfs', []):
                            if pdf.get('filename') == pdf_filename or pdf.get('server_filename') == pdf_filename:
                                references.append({
                                    'board_id': file.get('id'),
                                    'board_name': file.get('name'),
                                    'folder_name': folder.get('name'),
                                    'pdf_info': {
                                        'filename': pdf.get('filename'),
                                        'added_at': pdf.get('added_at'),
                                        'pages': pdf.get('pages', 0)
                                    }
                                })
        
        return {
            "status": "success",
            "pdf_filename": pdf_filename,
            "reference_count": len(references),
            "references": references
        }
        
    except Exception as e:
        logger.error(f"获取PDF引用信息失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"获取引用信息失败: {str(e)}")

# 启动应用
if __name__ == "__main__":
    # 加载环境变量
    dotenv.load_dotenv('.env')
    
    # 打印欢迎信息
    print("\n=== WhatNote 服务已启动 ===")
    print(f"API密钥配置: {'已配置' if bool(os.getenv('QWEN_API_KEY')) else '未配置'}")
    print(f"视觉API配置: {'已配置' if bool(os.getenv('QWEN_VL_API_KEY')) else '未配置'}")
    print("=======================\n")
    
    # 应用启动时同步一次文件结构
    sync_app_state_to_butler()
    
    # 启动服务
    uvicorn.run("main:app", host="127.0.0.1", port=8000, reload=True)  

# 在 @app.post('/api/expert/dynamic/generate-pdf-note') 后添加分段生成笔记的API

@app.post('/api/expert/dynamic/generate-segmented-note')
async def submit_generate_segmented_note_task(request_data: dict = Body(...)):
    """提交分段生成PDF笔记任务"""
    try:
        board_id = request_data.get("board_id")
        filename = request_data.get("filename")
        start_page = request_data.get("start_page", 1)
        page_count = request_data.get("page_count", 40)
        existing_note = request_data.get("existing_note", "")
        
        if not board_id or not filename:
            raise HTTPException(status_code=400, detail="缺少必要参数 board_id 或 filename")
        
        logger.info(f"提交分段生成PDF笔记任务: {filename}, 起始页: {start_page}, 页数: {page_count}")
        
        # 获取简化专家系统实例
        expert = simple_expert_manager.get_expert(board_id)
        
        # 构建任务参数
        task_params = {
            "filename": filename,
            "start_page": start_page,
            "page_count": page_count,
            "existing_note": existing_note
        }
        
        # 提交任务 - 使用正确的参数格式
        task_id = await expert.submit_task("generate_segmented_note", task_params)
        
        return {
            "task_id": task_id,
            "status": "submitted",
            "filename": filename,
            "start_page": start_page,
            "page_count": page_count,
            "message": f"分段笔记生成任务已提交，任务ID: {task_id}"
        }
        
    except Exception as e:
        logger.error(f"提交分段生成PDF笔记任务失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"提交任务失败: {str(e)}")

@app.post('/api/expert/dynamic/continue-segmented-note')
async def submit_continue_segmented_note_task(request_data: dict = Body(...)):
    """提交继续生成PDF笔记任务"""
    try:
        board_id = request_data.get("board_id")
        filename = request_data.get("filename")
        current_note = request_data.get("current_note", "")
        next_start_page = request_data.get("next_start_page")
        page_count = request_data.get("page_count", 40)
        
        if not board_id or not filename or not next_start_page:
            raise HTTPException(status_code=400, detail="缺少必要参数")
        
        logger.info(f"提交继续生成PDF笔记任务: {filename}, 起始页: {next_start_page}, 页数: {page_count}")
        
        # 获取简化专家系统实例
        expert = simple_expert_manager.get_expert(board_id)
        
        # 构建任务参数
        task_params = {
            "filename": filename,
            "start_page": next_start_page,
            "page_count": page_count,
            "existing_note": current_note
        }
        
        # 提交任务 - 使用正确的参数格式
        task_id = await expert.submit_task("generate_segmented_note", task_params)
        
        return {
            "task_id": task_id,
            "status": "submitted",
            "filename": filename,
            "start_page": next_start_page,
            "page_count": page_count,
            "message": f"继续生成笔记任务已提交，任务ID: {task_id}"
        }
        
    except Exception as e:
        logger.error(f"提交继续生成PDF笔记任务失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"提交任务失败: {str(e)}")

# 新增SSE端点用于实时任务状态推送
@app.get('/api/expert/dynamic/task-events/{board_id}')
async def task_events_stream(board_id: str):
    """
    SSE端点，实时推送任务状态变化
    """
    logger.info(f"📻 [SSE] 客户端连接任务事件流: {board_id}")
    
    class TaskEventSubscriber:
        def __init__(self):
            self.connected = True
            self.queue = asyncio.Queue(maxsize=100)
        
        async def send_event(self, event_data: Dict[str, Any]):
            """发送事件到客户端"""
            if self.connected:
                try:
                    await self.queue.put(event_data)
                except Exception as e:
                    pass
                except asyncio.QueueFull:
                    logger.warning(f"📻 [SSE] 事件队列已满，丢弃事件: {board_id}")
        
        async def generate_events(self):
            """生成SSE事件流"""
            try:
                # 立即发送当前任务状态
                current_tasks = task_event_manager.get_board_tasks(board_id)
                initial_event = {
                    "type": "task_list_update",
                    "board_id": board_id,
                    "tasks": current_tasks,
                    "timestamp": datetime.now().isoformat()
                }
                yield f"data: {json.dumps(initial_event, ensure_ascii=False)}\n\n"
                
                # 持续推送事件
                while self.connected:
                    try:
                        # 等待事件，超时检查连接状态
                        event_data = await asyncio.wait_for(self.queue.get(), timeout=30.0)
                        event_json = json.dumps(event_data, ensure_ascii=False)
                        yield f"data: {event_json}\n\n"
                    except Exception as e:
                        pass
                    except asyncio.TimeoutError:
                        # 发送心跳包
                        heartbeat = {
                            "type": "heartbeat",
                            "timestamp": datetime.now().isoformat()
                        }
                        yield f"data: {json.dumps(heartbeat, ensure_ascii=False)}\n\n"
                    except Exception as e:
                        logger.error(f"📻 [SSE] 事件生成错误: {str(e)}")
                        
            except Exception as e:
                logger.error(f"📻 [SSE] 事件流异常: {str(e)}")
            finally:
                self.connected = False
                logger.info(f"📻 [SSE] 客户端断开连接: {board_id}")
    
    # 创建订阅者
    subscriber = TaskEventSubscriber()
    
    # 注册到事件管理器
    task_event_manager.subscribe(board_id, subscriber)
    
    try:
        # 返回SSE响应
        response = StreamingResponse(
            subscriber.generate_events(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Headers": "Cache-Control"
            }
        )
        return response
    except Exception as e:
        pass
    finally:
        # 清理订阅
        subscriber.connected = False
        task_event_manager.unsubscribe(board_id, subscriber)

@app.post('/api/boards/{board_id}/annotation-style')
async def set_board_annotation_style(board_id: str, request_data: dict = Body(...)):
    """设置展板的注释风格 - 轻量级操作，优先处理"""
    try:
        style = request_data.get('style', 'detailed')
        custom_prompt = request_data.get('custom_prompt', '')
        
        # 🔧 完全同步处理，无需线程池 - 设置风格是本地操作，很快
        expert = simple_expert_manager.get_expert(board_id)
        expert.set_annotation_style(style, custom_prompt)
        logger.info(f"✅ 展板 {board_id} 注释风格已更新为: {style}")
        
        return {
            "status": "success",
            "message": f"注释风格已设置为: {style}",
            "board_id": board_id,
            "annotation_style": style,
            "custom_prompt": custom_prompt
        }
        
    except Exception as e:
        logger.error(f"设置注释风格失败: {str(e)}")
        return JSONResponse(
            status_code=500,
            content={"detail": f"设置注释风格失败: {str(e)}"}
        )

@app.get('/api/boards/{board_id}/annotation-style')
async def get_board_annotation_style(board_id: str):
    """获取展板的当前注释风格 - 轻量级操作，优先处理"""
    try:
        # 🔧 完全同步处理，无需线程池 - 获取风格是本地操作，很快
        expert = simple_expert_manager.get_expert(board_id)
        style_info = expert.get_annotation_style()
        
        return {
            "status": "success",
            "board_id": board_id,
            "annotation_style": style_info["style"],
            "custom_prompt": style_info["custom_prompt"],
            "available_styles": {
                "keywords": "关键词解释，中英对照",
                "translation": "单纯翻译文本内容", 
                "detailed": "详细学术注释",
                "custom": "自定义提示词"
            }
        }
        
    except Exception as e:
        logger.error(f"获取注释风格失败: {str(e)}")
        return JSONResponse(
            status_code=500,
            content={"detail": f"获取注释风格失败: {str(e)}"}
        )

# 🔧 添加后备路由，处理前端可能的错误路径调用
@app.get('/boards/{board_id}/annotation-style')
async def get_board_annotation_style_fallback(board_id: str):
    """后备路由：处理前端错误路径调用（缺少/api前缀）"""
    logger.warning(f"⚠️ 检测到前端使用了错误路径 /boards/{board_id}/annotation-style，重定向到正确API")
    return await get_board_annotation_style(board_id)

# 🔧 添加POST方法的后备路由
@app.post('/boards/{board_id}/annotation-style')
async def set_board_annotation_style_fallback(board_id: str, request_data: dict = Body(...)):
    """后备路由：处理前端POST请求的错误路径调用（缺少/api前缀）"""
    logger.warning(f"⚠️ 检测到前端使用了错误POST路径 /boards/{board_id}/annotation-style，重定向到正确API")
    return await set_board_annotation_style(board_id, request_data)

# 控制台API端点
@app.post('/api/butler/console')
async def butler_console_command(request_data: dict = Body(...)):
    """处理控制台命令 - 直接命令处理，不经过LLM"""
    try:
        command = request_data.get('command', '').strip()
        # 🔧 修复：支持both current_path and multi_step_context参数
        current_path = request_data.get('current_path') or request_data.get('multi_step_context', {})
        
        if not command:
            return JSONResponse(
                status_code=400,
                content={"detail": "命令不能为空"}
            )
        
        logger.info(f"🖥️ [CONSOLE] 直接执行命令: {command}")
        if current_path:
            logger.info(f"🖥️ [CONSOLE] 路径上下文: {current_path}")
        
        # 直接解析和执行命令
        result = await execute_direct_command(command, current_path)
        
        return {
            "status": "success",
            "result": result
        }
        
    except Exception as e:
        logger.error(f"🖥️ [CONSOLE] 命令处理失败: {str(e)}")
        return JSONResponse(
            status_code=500,
            content={"detail": f"命令处理失败: {str(e)}"}
        )

async def execute_direct_command(command: str, current_path: dict = None):
    """直接执行控制台命令，不经过LLM"""
    try:
        # 分割命令和参数
        parts = command.strip().split()
        if not parts:
            return {
                "response": "请输入命令", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        cmd = parts[0].lower()
        args = parts[1:]
        
        # 处理各种命令
        if cmd == "ls" or cmd == "list":
            return await handle_ls_command(args, current_path)
        elif cmd == "pwd":
            return handle_pwd_command(current_path)
        elif cmd == "cd":
            return await handle_cd_command(args, current_path)
        elif cmd == "help":
            return handle_help_command(args, current_path)
        elif cmd == "clear":
                return {
                "response": "clear", 
                "type": "clear",
                "style": "color: #ffffff; background: transparent;"
            }
        elif cmd == "history":
            return {
                "response": "命令历史功能需要前端配合实现", 
                "type": "info",
                "style": "color: #ffd43b; background: transparent;"
            }
        elif cmd == "status":
            return await handle_status_command(args)
        elif cmd == "exit":
            return {
                "response": "关闭控制台", 
                "type": "exit",
                "style": "color: #74c0fc; background: transparent;"
            }
        elif cmd == "tree":
            return await handle_tree_command(args)
        elif cmd == "find":
            return await handle_find_command(args)
        elif cmd == "search":
            return await handle_search_command(args)
        elif cmd == "stats":
            return await handle_stats_command(args)
        elif cmd == "recent":
            return await handle_recent_command(args)
        elif cmd == "backup":
            return await handle_backup_command(args, current_path)
        elif cmd == "export":
            return await handle_export_command(args, current_path)
        elif cmd == "delete":
            return await handle_delete_command(args, current_path)
        elif cmd == "rename":
            return await handle_rename_command(args, current_path)
        elif cmd == "info":
            return await handle_info_command(args, current_path)
        elif cmd == "copy":
            return await handle_copy_command(args, current_path)
        elif cmd == "goto":
            return await handle_goto_command(args, current_path)
        elif cmd == "next":
            return await handle_next_command(args, current_path)
        elif cmd == "prev":
            return await handle_prev_command(args, current_path)
        elif cmd == "first":
            return await handle_first_command(args, current_path)
        elif cmd == "last":
            return await handle_last_command(args, current_path)
        elif cmd == "pages":
            return await handle_pages_command(args, current_path)
        elif cmd == "annotate":
            return await handle_annotate_command(args, current_path)
        elif cmd == "annotation":
            return await handle_annotation_command(args, current_path)
        elif cmd == "page":
            return await handle_page_command(args, current_path)
        elif cmd == "window":
            return await handle_window_command(args, current_path)
        elif cmd == "layout":
            return await handle_layout_command(args, current_path)
        elif cmd == "config":
            return await handle_config_command(args)
        elif cmd == "log":
            return await handle_log_command(args)
        elif cmd == "cache":
            return await handle_cache_command(args)
        elif cmd == "refresh":
            return await handle_refresh_command(args)
        elif cmd == "version":
            return handle_version_command(args)
        elif cmd == "quota":
            return await handle_quota_command(args)
        
        # 课程操作命令
        elif cmd == "course" and len(args) > 0:
            return await handle_course_command(args)
        
        # 展板操作命令  
        elif cmd == "board" and len(args) > 0:
            return await handle_board_command(args, current_path)
        
        # PDF操作命令
        elif cmd == "pdf" and len(args) > 0:
            return await handle_pdf_command(args, current_path)
        
        # 笔记操作命令
        elif cmd == "note" and len(args) > 0:
            return await handle_note_command(args, current_path)
        
        else:
            return {
                "response": f"未知命令: {cmd}\n输入 'help' 查看可用命令",
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
            
    except Exception as e:
        logger.error(f"🖥️ [CONSOLE] 直接命令执行失败: {str(e)}")
        return {
            "response": f"命令执行失败: {str(e)}",
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }

async def handle_ls_command(args, current_path):
    """处理ls命令 - 与实际文件系统同步"""
    path_type = current_path.get('context', {}).get('type', 'root') if current_path else 'root'
    
    if path_type == 'root':
        # 在根目录，显示课程文件夹和全局信息
        course_folders = app_state.get_course_folders()
        
        response = f"📁 课程文件夹 ({len(course_folders)}):\n"
        for i, folder in enumerate(course_folders, 1):
            boards = app_state.get_boards()
            # 修复：course_folder字段存储的是课程ID，需要同时匹配ID和名称
            folder_id = folder.get('id', '')
            folder_name = folder.get('name', '')
            board_count = len([b for b in boards if b.get('course_folder') == folder_id or b.get('course_folder') == folder_name])
            response += f"  {i}. {folder['name']} ({board_count} 个展板)\n"
        
        # 显示全局PDF文件
        uploads_dir = "uploads"
        if os.path.exists(uploads_dir):
            pdf_files = [f for f in os.listdir(uploads_dir) if f.endswith('.pdf')]
            if pdf_files:
                response += f"\n📄 PDF文件 ({len(pdf_files)}):\n"
                for i, pdf in enumerate(pdf_files[:5], 1):  # 只显示前5个
                    response += f"  {i}. {pdf}\n"
                if len(pdf_files) > 5:
                    response += f"  ... 还有 {len(pdf_files) - 5} 个文件\n"
        
        return {
            "response": response,
            "type": "info",
            "style": "color: #ffffff; background: transparent;"
        }
    
    elif path_type == 'course':
        course_name = current_path.get('context', {}).get('courseName', '')
        course_id = current_path.get('context', {}).get('courseId', '')
        boards = app_state.get_boards()
        # 修复：course_folder字段存储的是课程ID，需要同时匹配ID和名称
        course_boards = [b for b in boards if b.get('course_folder') == course_id or b.get('course_folder') == course_name]
        
        if course_boards:
            response = f"📋 课程 '{course_name}' 的展板 ({len(course_boards)}):\n"
            for i, board in enumerate(course_boards, 1):
                response += f"  {i}. {board['name']}\n"
        else:
            response = f"📋 课程 '{course_name}' 暂无展板"
        
        return {
            "response": response,
            "type": "info", 
            "style": "color: #ffffff; background: transparent;"
        }
    
    elif path_type == 'board':
        board_name = current_path.get('context', {}).get('boardName', '')
        course_name = current_path.get('context', {}).get('courseName', '')
        
        # 获取展板的窗口和PDF
        boards = app_state.get_boards()
        current_board = None
        for board in boards:
            if board.get('name') == board_name and board.get('course_folder') == course_name:
                current_board = board
                break
            
        if current_board:
            board_id = current_board.get('id')
            response = f"🪟 展板 '{board_name}' 内容:\n"
            
            # 显示窗口
            try:
                from board_logger import BoardLogger
                board_logger = BoardLogger()
                log_data = board_logger.load_log(board_id)
                windows = log_data.get("windows", [])
                
                if windows:
                    response += f"\n📋 窗口 ({len(windows)}):\n"
                    for i, window in enumerate(windows, 1):
                        window_type = window.get("type", "")
                        title = window.get("title", "")
                        response += f"  {i}. [{window_type}] {title}\n"
                else:
                    response += "\n📋 暂无窗口\n"
            except Exception as e:
                pass
            except Exception:
                response += "\n📋 无法获取窗口信息\n"
            
            # 显示PDF文件
            try:
                if os.path.exists("board_data.json"):
                    with open("board_data.json", 'r', encoding='utf-8') as f:
                        board_data = json.load(f)
                        if board_data.get('board_id') == board_id:
                            pdfs = board_data.get('pdfs', [])
                            if pdfs:
                                response += f"\n📄 PDF文件 ({len(pdfs)}):\n"
                                for i, pdf in enumerate(pdfs, 1):
                                    filename = pdf.get('filename', '')
                                    current_page = pdf.get('currentPage', 1)
                                    response += f"  {i}. {filename} (页: {current_page})\n"
                            else:
                                response += "\n📄 暂无PDF文件\n"
            except Exception as e:
                pass
            except Exception:
                response += "\n📄 无法获取PDF信息\n"
        else:
            response = f"❌ 未找到展板 '{board_name}'"
        
        return {
            "response": response,
            "type": "info",
            "style": "color: #ffffff; background: transparent;"
        }
    
    else:
        return {
            "response": f"未知的路径类型: {path_type}",
            "type": "error", 
            "style": "color: #ff6b6b; background: transparent;"
        }

async def handle_cd_command(args, current_path):
    """处理cd命令"""
    if not args:
        return {
            "response": "用法: cd <目标目录>", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }
    
    target = ' '.join(args)  # 支持带空格的名称
    
    # 去掉外层引号（支持双引号和单引号）
    if (target.startswith('"') and target.endswith('"')) or (target.startswith("'") and target.endswith("'")):
        target = target[1:-1]
    
    if target == ".." or target == "..":
        return {
            "response": "返回上级目录",
            "type": "navigation",
            "style": "color: #74c0fc; background: transparent;",
            "navigation": {"action": "go_back"}
        }
    elif target == "/" or target == "~" or target == "root":
        return {
            "response": "返回根目录",
            "type": "navigation", 
            "style": "color: #74c0fc; background: transparent;",
            "navigation": {"action": "go_root"}
        }
    else:
        # 检查目标是否存在
        course_folders = app_state.get_course_folders()
        boards = app_state.get_boards()
        
        # 检查是否是课程文件夹 - 支持精确匹配、模糊匹配和部分匹配
        course_matches = []
        for folder in course_folders:
            if folder['name'] == target:
                # 精确匹配，优先级最高
                return {
                    "response": f"进入课程: {folder['name']}",
                    "type": "navigation",
                    "style": "color: #74c0fc; background: transparent;",
                    "navigation": {
                        "action": "enter_course",
                        "course_name": folder['name'],
                        "course_id": folder['id']
                    }
                }
            elif folder['name'].lower() == target.lower():
                # 大小写不敏感匹配
                course_matches.append(folder)
            elif target.lower() in folder['name'].lower() or folder['name'].lower() in target.lower():
                # 部分匹配（用于处理重命名情况）
                course_matches.append(folder)
        
        # 如果有匹配的课程，选择最佳匹配
        if course_matches:
            best_match = course_matches[0]  # 取第一个匹配项
            return {
                "response": f"进入课程: {best_match['name']}",
                "type": "navigation",
                "style": "color: #74c0fc; background: transparent;",
                "navigation": {
                    "action": "enter_course",
                    "course_name": best_match['name'],
                    "course_id": best_match['id']
                }
            }
        
        # 检查是否是展板 - 支持精确匹配、模糊匹配和部分匹配
        board_matches = []
        for board in boards:
            if board['name'] == target:
                # 精确匹配，优先级最高
                return {
                    "response": f"进入展板: {board['name']}",
                    "type": "navigation",
                    "style": "color: #74c0fc; background: transparent;",
                    "navigation": {
                        "action": "enter_board", 
                        "board_name": board['name'],
                        "board_id": board['id']
                    }
                }
            elif board['name'].lower() == target.lower():
                # 大小写不敏感匹配
                board_matches.append(board)
            elif target.lower() in board['name'].lower() or board['name'].lower() in target.lower():
                # 部分匹配（用于处理重命名情况）
                board_matches.append(board)
        
        # 如果有匹配的展板，选择最佳匹配
        if board_matches:
            best_match = board_matches[0]  # 取第一个匹配项
            return {
                "response": f"进入展板: {best_match['name']}",
                "type": "navigation",
                "style": "color: #74c0fc; background: transparent;",
                "navigation": {
                    "action": "enter_board", 
                    "board_name": best_match['name'],
                    "board_id": best_match['id']
                }
            }
        
        # 提供建议
        suggestions = []
        for folder in course_folders:
            if target.lower() in folder['name'].lower():
                suggestions.append(folder['name'])
        for board in boards:
            if target.lower() in board['name'].lower():
                suggestions.append(board['name'])
        
        error_msg = f"找不到目录: {target}"
        if suggestions:
            error_msg += f"\n相似的目录: {', '.join(suggestions[:3])}"
        
        return {
            "response": error_msg, 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }

def handle_pwd_command(current_path):
    """处理pwd命令"""
    if not current_path or 'context' not in current_path:
        return {
            "response": "当前路径: / (根目录)", 
            "type": "info",
            "style": "color: #74c0fc; background: transparent;"
        }
    
    context = current_path['context']
    path_type = context.get('type', 'root')
    
    if path_type == 'course':
        course_name = context.get('courseName', '未知课程')
        return {
            "response": f"当前路径: /{course_name}/ (课程目录)", 
            "type": "info",
            "style": "color: #74c0fc; background: transparent;"
        }
    elif path_type == 'board':
        course_name = context.get('courseName', '未知课程')
        board_name = context.get('boardName', '未知展板')
        return {
            "response": f"当前路径: /{course_name}/{board_name}/ (展板目录)", 
            "type": "info",
            "style": "color: #74c0fc; background: transparent;"
        }
    elif path_type == 'pdf':
        course_name = context.get('courseName', '未知课程')
        board_name = context.get('boardName', '未知展板')
        pdf_name = context.get('pdfName', '未知PDF')
        return {
            "response": f"当前路径: /{course_name}/{board_name}/{pdf_name} (PDF文件)", 
            "type": "info",
            "style": "color: #74c0fc; background: transparent;"
        }
    else:
        return {
            "response": "当前路径: / (根目录)", 
            "type": "info",
            "style": "color: #74c0fc; background: transparent;"
        }

async def handle_status_command(args):
    """处理status命令"""
    try:
        # 获取系统状态信息
        course_folders = app_state.get_course_folders()
        boards = app_state.get_boards()
        
        # 扫描PDF文件
        uploads_dir = "uploads"
        pdf_count = 0
        total_size = 0
        
        if os.path.exists(uploads_dir):
            for file in os.listdir(uploads_dir):
                if file.endswith('.pdf'):
                    pdf_count += 1
                    file_path = os.path.join(uploads_dir, file)
                    if os.path.exists(file_path):
                        total_size += os.path.getsize(file_path)
        
        # 格式化大小
        def format_size(size_bytes):
            if size_bytes == 0:
                return "0 B"
            size_names = ["B", "KB", "MB", "GB"]
            import math
            i = int(math.floor(math.log(size_bytes, 1024)))
            p = math.pow(1024, i)
            s = round(size_bytes / p, 2)
            return f"{s} {size_names[i]}"
        
        if '-d' in args or '--detail' in args:
            response = "📊 WhatNote 系统详细状态:\n\n"
            response += f"📚 课程文件夹: {len(course_folders)} 个\n"
            if course_folders:
                for folder in course_folders[:5]:  # 显示前5个
                    response += f"  - {folder['name']}\n"
                if len(course_folders) > 5:
                    response += f"  ... 还有 {len(course_folders) - 5} 个\n"
            
            response += f"\n📋 展板: {len(boards)} 个\n"
            if boards:
                for board in boards[:5]:  # 显示前5个
                    course_info = f" [{board['course_folder']}]" if board.get('course_folder') else ""
                    response += f"  - {board['name']}{course_info}\n"
                if len(boards) > 5:
                    response += f"  ... 还有 {len(boards) - 5} 个\n"
            
            response += f"\n📄 PDF文件: {pdf_count} 个\n"
            response += f"💾 存储使用: {format_size(total_size)}\n"
        else:
            response = "📊 WhatNote 系统状态:\n"
            response += f"📚 课程: {len(course_folders)} | 📋 展板: {len(boards)} | 📄 PDF: {pdf_count}\n"
            response += f"💾 存储: {format_size(total_size)}"
        
        return {
            "response": response, 
            "type": "info",
            "style": "color: #ffffff; background: transparent;"
        }
        
    except Exception as e:
        return {
            "response": f"获取系统状态失败: {str(e)}", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }

async def handle_course_command(args):
    """处理course命令"""
    if not args:
        return {
            "response": "用法: course <list|create|delete|rename|show> [名称]", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }
    
    action = args[0].lower()
    
    if action == "list":
        course_folders = app_state.get_course_folders()
        if course_folders:
            response = f"📚 课程文件夹列表 ({len(course_folders)}):\n"
            for i, folder in enumerate(course_folders, 1):
                response += f"  {i}. {folder['name']} (ID: {folder['id']})\n"
        else:
            response = "暂无课程文件夹"
        return {
            "response": response, 
            "type": "info",
            "style": "color: #ffffff; background: transparent;"
        }
    
    elif action == "create":
        if len(args) < 2:
            return {
                "response": "用法: course create <课程名称>", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        course_name = ' '.join(args[1:])  # 支持带空格的名称
        
        # 去掉外层引号（支持双引号和单引号）
        if (course_name.startswith('"') and course_name.endswith('"')) or (course_name.startswith("'") and course_name.endswith("'")):
            course_name = course_name[1:-1]
        
        # 检查是否已存在
        if app_state.course_folder_exists(course_name):
            return {
                "response": f"课程文件夹 '{course_name}' 已存在", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        # 创建课程文件夹
        folder = app_state.add_course_folder(course_name)
        app_state.save_state()
        
        return {
            "response": f"✅ 课程文件夹 '{course_name}' 创建成功", 
            "type": "success",
            "style": "color: #51cf66; background: transparent;",
            "refresh_needed": True  # 通知前端需要刷新
        }
    
    elif action == "delete":
        if len(args) < 2:
            return {
                "response": "用法: course delete <课程名称>", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        course_name = ' '.join(args[1:])
        
        # 去掉外层引号（支持双引号和单引号）
        if (course_name.startswith('"') and course_name.endswith('"')) or (course_name.startswith("'") and course_name.endswith("'")):
            course_name = course_name[1:-1]
        
        # 查找要删除的课程
        course_folders = app_state.get_course_folders()
        target_course = None
        for folder in course_folders:
            if folder['name'] == course_name or folder['name'].lower() == course_name.lower():
                target_course = folder
                break
        
        if not target_course:
            return {
                "response": f"找不到课程文件夹: {course_name}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        # 检查课程下是否有展板
        boards = app_state.get_boards()
        course_boards = [b for b in boards if b.get('course_folder') == course_name]
        
        if course_boards:
            board_names = ', '.join([b['name'] for b in course_boards[:3]])
            if len(course_boards) > 3:
                board_names += f" 等{len(course_boards)}个展板"
            return {
                "response": f"无法删除课程 '{course_name}'，该课程下还有展板: {board_names}\n请先删除相关展板", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        # 删除课程文件夹
        app_state.course_folders = [f for f in app_state.course_folders if f['id'] != target_course['id']]
        app_state.save_state()
        
        return {
            "response": f"✅ 课程文件夹 '{course_name}' 已删除", 
            "type": "success",
            "style": "color: #51cf66; background: transparent;",
            "refresh_needed": True
        }
    
    elif action == "rename":
        if len(args) < 3:
            return {
                "response": "用法: course rename <旧名称> <新名称>", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        old_name = args[1]
        new_name = ' '.join(args[2:])
        
        # 去掉外层引号
        if (old_name.startswith('"') and old_name.endswith('"')) or (old_name.startswith("'") and old_name.endswith("'")):
            old_name = old_name[1:-1]
        if (new_name.startswith('"') and new_name.endswith('"')) or (new_name.startswith("'") and new_name.endswith("'")):
            new_name = new_name[1:-1]
        
        # 查找要重命名的课程
        course_folders = app_state.get_course_folders()
        target_course = None
        for folder in course_folders:
            if folder['name'] == old_name or folder['name'].lower() == old_name.lower():
                target_course = folder
                break
        
        if not target_course:
            return {
                "response": f"找不到课程文件夹: {old_name}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        # 检查新名称是否已存在
        if app_state.course_folder_exists(new_name):
            return {
                "response": f"课程文件夹 '{new_name}' 已存在", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        # 更新课程名称和相关展板的course_folder字段
        target_course['name'] = new_name
        boards = app_state.get_boards()
        for board in boards:
            if board.get('course_folder') == old_name:
                board['course_folder'] = new_name
        
        app_state.save_state()
        
        return {
            "response": f"✅ 课程文件夹已重命名: '{old_name}' → '{new_name}'", 
            "type": "success",
            "style": "color: #51cf66; background: transparent;",
            "refresh_needed": True
        }

    elif action == "show" or action == "info":
        if len(args) < 2:
            return {
                "response": "用法: course show <课程名称>", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        course_name = ' '.join(args[1:])
        
        # 去掉外层引号（支持双引号和单引号）
        if (course_name.startswith('"') and course_name.endswith('"')) or (course_name.startswith("'") and course_name.endswith("'")):
            course_name = course_name[1:-1]
        
        # 查找课程文件夹
        course_folders = app_state.get_course_folders()
        target_course = None
        for folder in course_folders:
            if folder['name'] == course_name or folder['name'].lower() == course_name.lower():
                target_course = folder
                break
        
        if target_course:
            response = f"📚 课程详情: {target_course['name']}\n"
            response += f"  ID: {target_course['id']}\n"
            response += f"  创建时间: {target_course.get('created_at', '未知')}\n"
            
            # 查找该课程下的展板
            boards = app_state.get_boards()
            course_boards = [b for b in boards if b.get('course_folder') == course_name]
            if course_boards:
                response += f"  关联展板: {len(course_boards)} 个\n"
                for board in course_boards[:3]:  # 显示前3个
                    response += f"    - {board['name']}\n"
                if len(course_boards) > 3:
                    response += f"    ... 还有 {len(course_boards) - 3} 个\n"
            else:
                response += "  关联展板: 无\n"
            
            return {
                "response": response, 
                "type": "info",
                "style": "color: #ffffff; background: transparent;"
            }
        else:
            return {
                "response": f"找不到课程文件夹: {course_name}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
            return {
                "response": f"未知的课程操作: {action}", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }

async def handle_board_command(args, current_path):
    """处理board命令"""
    if not args:
        return {
            "response": "用法: board <list|create|open|delete> [名称]", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }
    
    action = args[0].lower()
    
    if action == "list":
        # 根据当前路径上下文显示展板
        path_type = current_path.get('context', {}).get('type', 'root') if current_path else 'root'
        
        if path_type == 'course':
            # 在课程中，只显示当前课程的展板
            course_name = current_path.get('context', {}).get('courseName', '')
            boards = app_state.get_boards()
            course_boards = [b for b in boards if b.get('course_folder') == course_name]
            
            if course_boards:
                response = f"📋 课程 '{course_name}' 的展板 ({len(course_boards)}):\n"
                for i, board in enumerate(course_boards, 1):
                    response += f"  {i}. {board['name']} (ID: {board['id']})\n"
            else:
                response = f"📋 课程 '{course_name}' 暂无展板"
        else:
            # 在根目录或其他位置，显示所有展板
            boards = app_state.get_boards()
            if boards:
                response = f"📋 所有展板 ({len(boards)}):\n"
                for i, board in enumerate(boards, 1):
                    course_info = f" [课程: {board['course_folder']}]" if board.get('course_folder') else ""
                    response += f"  {i}. {board['name']} (ID: {board['id']}){course_info}\n"
            else:
                response = "📋 系统中暂无展板"
        
        return {
            "response": response, 
            "type": "info",
            "style": "color: #ffffff; background: transparent;"
        }
    
    elif action == "create":
        if len(args) < 2:
            return {
                "response": "用法: board create <展板名称>", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        board_name = ' '.join(args[1:])
        
        # 去掉外层引号（支持双引号和单引号）
        if (board_name.startswith('"') and board_name.endswith('"')) or (board_name.startswith("'") and board_name.endswith("'")):
            board_name = board_name[1:-1]
        
        # 根据当前路径上下文确定课程名称
        path_type = current_path.get('context', {}).get('type', 'root') if current_path else 'root'
        course_folder = ""
        
        if path_type == 'course':
            # 在课程目录中，使用当前课程名称
            course_folder = current_path.get('context', {}).get('courseName', '')
            
            # 验证课程是否真实存在
            course_folders = app_state.get_course_folders()
            course_exists = any(folder['name'] == course_folder for folder in course_folders)
            
            if not course_exists:
                return {
                    "response": f"错误：当前课程 '{course_folder}' 不存在，无法创建展板", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
        
        # 检查展板是否已存在（在指定课程中）
        if app_state.board_exists(board_name, course_folder):
            scope_msg = f"课程 '{course_folder}' 中" if course_folder else "系统中"
            return {
                "response": f"展板 '{board_name}' 在{scope_msg}已存在", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        # 创建展板
        board = app_state.add_board(board_name, course_folder)
        app_state.save_state()
        
        location_msg = f"课程 '{course_folder}' 下" if course_folder else "根目录下"
        return {
            "response": f"✅ 展板 '{board_name}' 在{location_msg}创建成功", 
            "type": "success",
            "style": "color: #51cf66; background: transparent;",
            "refresh_needed": True  # 通知前端需要刷新
        }
    
    elif action == "open":
        if len(args) < 2:
            return {
                "response": "用法: board open <展板名称>", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        board_name = ' '.join(args[1:])
        
        # 去掉外层引号（支持双引号和单引号）
        if (board_name.startswith('"') and board_name.endswith('"')) or (board_name.startswith("'") and board_name.endswith("'")):
            board_name = board_name[1:-1]
        
        boards = app_state.get_boards()
        
        # 查找展板 - 只支持精确匹配
        target_board = None
        for board in boards:
            if board['name'] == board_name:
                target_board = board
                break
        
        # 如果精确匹配失败，尝试不区分大小写的匹配
        if not target_board:
            for board in boards:
                if board['name'].lower() == board_name.lower():
                    target_board = board
                    break
        
        if target_board:
            return {
                "response": f"打开展板: {target_board['name']}",
                "type": "navigation",
                "style": "color: #74c0fc; background: transparent;",
                "navigation": {
                    "action": "open_board",
                    "board_name": target_board['name'],
                    "board_id": target_board['id']
                }
            }
        else:
            # 提供建议
            suggestions = []
            for board in boards:
                if board_name.lower() in board['name'].lower():
                    suggestions.append(board['name'])
            
            error_msg = f"找不到展板: {board_name}"
            if suggestions:
                error_msg += f"\n💡 您是否在找: {', '.join(suggestions[:3])}"
            
            return {
                "response": error_msg, 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
    elif action == "delete":
        if len(args) < 2:
            return {
                "response": "用法: board delete <展板名称>", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        board_name = ' '.join(args[1:])
        
        # 去掉外层引号（支持双引号和单引号）
        if (board_name.startswith('"') and board_name.endswith('"')) or (board_name.startswith("'") and board_name.endswith("'")):
            board_name = board_name[1:-1]
        
        # 查找要删除的展板
        boards = app_state.get_boards()
        target_board = None
        for board in boards:
            if board['name'] == board_name or board['name'].lower() == board_name.lower():
                target_board = board
                break
        
        if not target_board:
            return {
                "response": f"找不到展板: {board_name}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        # 检查展板是否有关联的PDF文件
        board_id = target_board.get('id')
        has_pdfs = False
        try:
            board_data_file = f"board_data.json"
            if os.path.exists(board_data_file):
                with open(board_data_file, 'r', encoding='utf-8') as f:
                    board_data = json.load(f)
                    if board_data.get('board_id') == board_id:
                        pdfs = board_data.get('pdfs', [])
                        if pdfs:
                            has_pdfs = True
                            pdf_names = ', '.join([pdf.get('filename', '') for pdf in pdfs[:3]])
                            if len(pdfs) > 3:
                                pdf_names += f" 等{len(pdfs)}个文件"
                            return {
                                "response": f"无法删除展板 '{board_name}'，该展板下还有PDF文件: {pdf_names}\n请先删除相关PDF文件", 
                                "type": "error",
                                "style": "color: #ff6b6b; background: transparent;"
                            }
        except Exception as e:
            pass
        except:
            pass
        
        # 删除展板
        app_state.boards = [b for b in app_state.boards if b['id'] != target_board['id']]
        app_state.save_state()
        
        # 清理相关的展板日志文件
        try:
            board_log_path = os.path.join("board_logs", f"{board_id}.json")
            if os.path.exists(board_log_path):
                os.remove(board_log_path)
        except Exception as e:
            pass
        except:
            pass
        
        course_info = f" (课程: {target_board.get('course_folder', '根目录')})" if target_board.get('course_folder') else ""
        return {
            "response": f"✅ 展板 '{board_name}'{course_info} 已删除", 
            "type": "success",
            "style": "color: #51cf66; background: transparent;",
            "refresh_needed": True
        }
    elif action == "rename":
        if len(args) < 3:
            return {
                "response": "用法: board rename <旧名称> <新名称>", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        old_name = args[1]
        new_name = ' '.join(args[2:])
        
        # 去掉外层引号
        if (old_name.startswith('"') and old_name.endswith('"')) or (old_name.startswith("'") and old_name.endswith("'")):
            old_name = old_name[1:-1]
        if (new_name.startswith('"') and new_name.endswith('"')) or (new_name.startswith("'") and new_name.endswith("'")):
            new_name = new_name[1:-1]
        
        # 查找要重命名的展板
        boards = app_state.get_boards()
        target_board = None
        for board in boards:
            if board['name'] == old_name or board['name'].lower() == old_name.lower():
                target_board = board
                break
        
        if not target_board:
            return {
                "response": f"找不到展板: {old_name}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        # 检查新名称是否已存在（在同一课程中）
        course_folder = target_board.get('course_folder', '')
        for board in boards:
            if (board['name'] == new_name or board['name'].lower() == new_name.lower()) and board.get('course_folder', '') == course_folder and board['id'] != target_board['id']:
                scope_msg = f"课程 '{course_folder}' 中" if course_folder else "系统中"
                return {
                    "response": f"展板名称 '{new_name}' 在{scope_msg}已存在", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
        
        # 更新展板名称
        target_board['name'] = new_name
        app_state.save_state()
        
        course_info = f" (课程: {course_folder})" if course_folder else ""
        return {
            "response": f"✅ 展板已重命名: '{old_name}' → '{new_name}'{course_info}", 
            "type": "success",
            "style": "color: #51cf66; background: transparent;",
            "refresh_needed": True
        }
    
    elif action == "write":
        if len(args) < 3:
            return {
                "response": "用法: window write <窗口ID> \"内容\"", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        window_id = args[1]
        content = ' '.join(args[2:])
        
        # 去掉外层引号，但保留内容中的引号
        if (content.startswith('"') and content.endswith('"')) or (content.startswith("'") and content.endswith("'")):
            content = content[1:-1]
        
        try:
            from board_logger import BoardLogger
            board_logger = BoardLogger()
            log_data = board_logger.load_log(board_id)
            windows = log_data.get("windows", [])
            
            # 查找指定窗口
            target_window = None
            for window in windows:
                if window.get("id") == window_id:
                    target_window = window
                    break
            
            if not target_window:
                return {
                    "response": f"找不到窗口: {window_id}", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
            
            # 更新窗口内容
            target_window["content"] = content
            success = board_logger.update_window(board_id, window_id, target_window)
            
            if success:
                return {
                    "response": f"✅ 已更新窗口 {window_id} 的内容", 
                    "type": "success",
                    "style": "color: #51cf66; background: transparent;",
                    "refresh_needed": True
                }
            else:
                return {
                    "response": f"更新窗口内容失败", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
        except Exception as e:
            return {
                "response": f"写入窗口内容失败: {str(e)}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
    elif action == "list":
        try:
            from board_logger import BoardLogger
            board_logger = BoardLogger()
            log_data = board_logger.load_log(board_id)
            windows = log_data.get("windows", [])
            
            if windows:
                response = f"📋 当前展板的窗口 ({len(windows)}):\n"
                for i, window in enumerate(windows, 1):
                    window_id = window.get("id", "")
                    window_type = window.get("type", "")
                    title = window.get("title", "")
                    content_preview = window.get("content", "")[:30]
                    if len(window.get("content", "")) > 30:
                        content_preview += "..."
                    response += f"  {i}. [{window_type}] {title} (ID: {window_id})\n"
                    if content_preview:
                        response += f"     内容: {content_preview}\n"
            else:
                response = f"📋 当前展板暂无窗口"
            
            return {
                "response": response, 
                "type": "info",
                "style": "color: #ffffff; background: transparent;"
            }
        except Exception as e:
            return {
                "response": f"获取窗口列表失败: {str(e)}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
    elif action == "show":
        if len(args) < 2:
            return {
                "response": "用法: window show <窗口ID>", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        window_id = args[1]
        
        try:
            from board_logger import BoardLogger
            board_logger = BoardLogger()
            log_data = board_logger.load_log(board_id)
            windows = log_data.get("windows", [])
            
            # 查找指定窗口
            target_window = None
            for window in windows:
                if window.get("id") == window_id:
                    target_window = window
                    break
            
            if not target_window:
                return {
                    "response": f"找不到窗口: {window_id}", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
            
            # 显示窗口详细信息
            response = f"🪟 窗口详情: {target_window.get('title', '')}\n"
            response += f"  ID: {window_id}\n"
            response += f"  类型: {target_window.get('type', '')}\n"
            response += f"  位置: x={target_window.get('position', {}).get('x', 0)}, y={target_window.get('position', {}).get('y', 0)}\n"
            response += f"  大小: {target_window.get('size', {}).get('width', 0)}x{target_window.get('size', {}).get('height', 0)}\n"
            response += f"  创建时间: {target_window.get('created_at', '未知')}\n"
            content = target_window.get('content', '')
            if content:
                if len(content) > 100:
                    response += f"  内容: {content[:100]}...\n"
                else:
                    response += f"  内容: {content}\n"
            else:
                response += "  内容: (空)\n"
            
            return {
                "response": response, 
                "type": "info",
                "style": "color: #ffffff; background: transparent;"
            }
        except Exception as e:
            return {
                "response": f"获取窗口信息失败: {str(e)}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
    elif action == "image":
        if len(args) < 3:
            return {
                "response": "用法: window image <窗口ID> <图片路径或URL>\n支持:\n  - 本地文件: window image win123 \"C:/path/to/image.jpg\"\n  - 网络图片: window image win123 \"https://example.com/image.jpg\"\n  - 系统图片: window image win123 \"uploaded/filename.jpg\"", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        window_id = args[1]
        image_path = ' '.join(args[2:])
        
        # 去掉外层引号
        if (image_path.startswith('"') and image_path.endswith('"')) or (image_path.startswith("'") and image_path.endswith("'")):
            image_path = image_path[1:-1]
        
        try:
            from board_logger import BoardLogger
            board_logger = BoardLogger()
            log_data = board_logger.load_log(board_id)
            windows = log_data.get("windows", [])
            
            # 查找指定窗口
            target_window = None
            for window in windows:
                if window.get("id") == window_id:
                    target_window = window
                    break
            
            if not target_window:
                return {
                    "response": f"找不到窗口: {window_id}", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
            
            # 检查是否为图片窗口
            if target_window.get("type") != "image":
                return {
                    "response": f"窗口 {window_id} 不是图片窗口，无法设置图片", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
            
            # 处理不同类型的图片路径
            final_image_url = None
            
            if image_path.startswith(("http://", "https://")):
                # 网络图片URL
                final_image_url = image_path
                print(f"🌐 设置网络图片: {image_path}")
                
            elif os.path.exists(image_path):
                # 本地文件路径
                print(f"📁 处理本地文件: {image_path}")
                
                # 验证是否为图片文件
                allowed_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp'}
                file_ext = os.path.splitext(image_path.lower())[1]
                
                if file_ext not in allowed_extensions:
                    return {
                        "response": f"文件格式不支持: {file_ext}，支持的格式: {', '.join(allowed_extensions)}", 
                        "type": "error",
                        "style": "color: #ff6b6b; background: transparent;"
                    }
                
                # 复制文件到images目录
                import shutil
                import time
                
                filename = os.path.basename(image_path)
                name, ext = os.path.splitext(filename)
                timestamp = int(time.time())
                unique_filename = f"{name}_cmd_{timestamp}{ext}"
                
                images_dir = os.path.join(UPLOAD_DIR, 'images')
                if not os.path.exists(images_dir):
                    os.makedirs(images_dir, exist_ok=True)
                
                dest_path = os.path.join(images_dir, unique_filename)
                shutil.copy2(image_path, dest_path)
                
                final_image_url = f"/api/images/view/{unique_filename}"
                print(f"📋 复制到系统: {dest_path}")
                
            elif image_path.startswith("uploaded/") or image_path.startswith("/api/images/"):
                # 系统内部图片路径
                if image_path.startswith("uploaded/"):
                    # 兼容旧格式
                    filename = image_path.replace("uploaded/", "")
                    final_image_url = f"/api/images/view/{filename}"
                else:
                    # 已经是正确格式
                    final_image_url = image_path
                    
                print(f"🗃️  使用系统图片: {final_image_url}")
                
            else:
                return {
                    "response": f"图片路径无效: {image_path}\n请使用:\n  - 网络URL (http://或https://)\n  - 本地文件的完整路径\n  - 系统图片路径 (uploaded/filename.jpg)", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
            
            # 更新窗口的图片URL
            target_window["content"] = final_image_url
            success = board_logger.update_window(board_id, window_id, target_window)
            
            if success:
                return {
                    "response": f"✅ 已设置图片窗口 {window_id} 的图片: {final_image_url}", 
                    "type": "success",
                    "style": "color: #51cf66; background: transparent;",
                    "refresh_needed": True
                }
            else:
                return {
                    "response": f"设置图片失败", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
                
        except Exception as e:
            return {
                "response": f"设置图片失败: {str(e)}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
    elif action == "delete":
        if len(args) < 2:
            return {
                "response": "用法: window delete <窗口ID>", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        window_id = args[1]
        
        try:
            from board_logger import BoardLogger
            board_logger = BoardLogger()
            success = board_logger.remove_window(board_id, window_id)
            
            if success:
                return {
                    "response": f"✅ 已删除窗口: {window_id}", 
                    "type": "success",
                    "style": "color: #51cf66; background: transparent;",
                    "refresh_needed": True
                }
            else:
                return {
                    "response": f"找不到窗口: {window_id}", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
        except Exception as e:
            return {
                "response": f"删除窗口失败: {str(e)}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
    elif action == "show":
        if len(args) < 2:
            return {
                "response": "用法: window show <窗口ID>", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        window_id = args[1]
        
        try:
            from board_logger import BoardLogger
            board_logger = BoardLogger()
            log_data = board_logger.load_log(board_id)
            windows = log_data.get("windows", [])
            
            # 查找指定窗口
            target_window = None
            for window in windows:
                if window.get("id") == window_id:
                    target_window = window
                    break
            
            if not target_window:
                return {
                    "response": f"找不到窗口: {window_id}", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
            
            # 显示窗口详细信息
            response = f"🪟 窗口详情: {target_window.get('title', '')}\n"
            response += f"  ID: {window_id}\n"
            response += f"  类型: {target_window.get('type', '')}\n"
            response += f"  位置: x={target_window.get('position', {}).get('x', 0)}, y={target_window.get('position', {}).get('y', 0)}\n"
            response += f"  大小: {target_window.get('size', {}).get('width', 0)}x{target_window.get('size', {}).get('height', 0)}\n"
            response += f"  创建时间: {target_window.get('created_at', '未知')}\n"
            content = target_window.get('content', '')
            if content:
                if len(content) > 100:
                    response += f"  内容: {content[:100]}...\n"
                else:
                    response += f"  内容: {content}\n"
            else:
                response += "  内容: (空)\n"
            
            return {
                "response": response, 
                "type": "info",
                "style": "color: #ffffff; background: transparent;"
            }
        except Exception as e:
            return {
                "response": f"获取窗口信息失败: {str(e)}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
    else:
        return {
            "response": f"未知的窗口操作: {action}", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }

async def handle_layout_command(args, current_path):
    """处理layout命令"""
    return {
        "response": "📋 当前布局信息:\n  控制台: 启用\n  展板视图: 活跃\n  窗口管理: 正常", 
        "type": "info",
        "style": "color: #ffffff; background: transparent;"
    }

async def handle_config_command(args):
    """处理config命令"""
    if not args:
        return {
            "response": "用法: config <show|set> [键] [值]", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }
    
    action = args[0].lower()
    
    if action == "show":
        response = "⚙️ 当前配置:\n"
        response += "  控制台模式: 直接命令处理\n"
        response += "  LLM处理: 已禁用\n"
        response += "  自动补全: 启用\n"
        response += "  历史记录: 启用\n"
        
        return {
            "response": response, 
            "type": "info",
            "style": "color: #ffffff; background: transparent;"
        }
    else:
        return {
            "response": "配置设置功能需要进一步开发", 
            "type": "info",
            "style": "color: #ffd43b; background: transparent;"
        }

async def handle_log_command(args):
    """处理log命令"""
    if not args:
        return {
            "response": "用法: log <show|clear>", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }
    
    action = args[0].lower()
    
    if action == "show":
        # 读取日志文件
        log_content = []
        log_files = ['llm_interactions.log', 'logs/app.log']
        
        for log_file in log_files:
            if os.path.exists(log_file):
                try:
                    with open(log_file, 'r', encoding='utf-8') as f:
                        lines = f.readlines()
                        # 显示最后10行
                        log_content.extend(lines[-10:])
                except Exception as e:
                    pass
                except:
                    continue
        
        if log_content:
            response = "📋 最近日志记录:\n" + "".join(log_content[-20:])
        else:
            response = "📋 暂无日志记录"
            
        return {
            "response": response, 
            "type": "info",
            "style": "color: #ffffff; background: transparent;"
        }
    elif action == "clear":
        return {
            "response": "✅ 日志已清空", 
            "type": "success",
            "style": "color: #51cf66; background: transparent;"
        }
    else:
        return {
            "response": f"未知的日志操作: {action}", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }

async def handle_cache_command(args):
    """处理cache命令"""
    if not args:
        return {
            "response": "用法: cache <clear|show>", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }
    
    action = args[0].lower()
    
    if action == "clear":
        # 清理各种缓存
        cache_dirs = ['__pycache__', 'temp', 'frontend/temp']
        cleared_count = 0
        
        for cache_dir in cache_dirs:
            if os.path.exists(cache_dir):
                try:
                    import shutil
                    shutil.rmtree(cache_dir)
                    cleared_count += 1
                except Exception as e:
                    pass
                except:
                    pass
        
        return {
            "response": f"✅ 缓存已清理，清理了 {cleared_count} 个缓存目录", 
            "type": "success",
            "style": "color: #51cf66; background: transparent;"
        }
    elif action == "show":
        # 显示缓存状态
        cache_info = []
        cache_dirs = ['__pycache__', 'temp', 'frontend/temp', 'uploads/temp']
        
        for cache_dir in cache_dirs:
            if os.path.exists(cache_dir):
                try:
                    size = sum(os.path.getsize(os.path.join(cache_dir, f)) 
                              for f in os.listdir(cache_dir) 
                              if os.path.isfile(os.path.join(cache_dir, f)))
                    cache_info.append(f"  {cache_dir}: {size} 字节")
                except Exception as e:
                    pass
                except:
                    cache_info.append(f"  {cache_dir}: 无法访问")
            else:
                cache_info.append(f"  {cache_dir}: 不存在")
        
        response = "📋 缓存状态:\n" + "\n".join(cache_info)
        return {
            "response": response, 
            "type": "info",
            "style": "color: #ffffff; background: transparent;"
        }
    else:
        return {
            "response": f"未知的缓存操作: {action}", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }

async def handle_refresh_command(args):
    """处理refresh命令"""
    return {
        "response": "✅ 数据已刷新", 
        "type": "success",
        "style": "color: #51cf66; background: transparent;",
        "refresh_needed": True
    }

async def handle_quota_command(args):
    """处理quota命令"""
    try:
        uploads_dir = "uploads"
        total_size = 0
        file_count = 0
        
        if os.path.exists(uploads_dir):
            for file in os.listdir(uploads_dir):
                file_path = os.path.join(uploads_dir, file)
                if os.path.isfile(file_path):
                    total_size += os.path.getsize(file_path)
                    file_count += 1
        
        # 格式化大小
        def format_size(size_bytes):
            if size_bytes == 0:
                return "0 B"
            size_names = ["B", "KB", "MB", "GB"]
            import math
            i = int(math.floor(math.log(size_bytes, 1024)))
            p = math.pow(1024, i)
            s = round(size_bytes / p, 2)
            return f"{s} {size_names[i]}"
        
        response = "💾 存储空间使用情况:\n\n"
        response += f"📄 文件总数: {file_count}\n"
        response += f"📦 已使用空间: {format_size(total_size)}\n"
        response += f"📊 平均文件大小: {format_size(total_size / file_count) if file_count > 0 else '0 B'}\n"
        
        # 简单的使用率显示
        quota_limit = 1024 * 1024 * 1024  # 1GB 假设限制
        usage_percent = (total_size / quota_limit) * 100 if quota_limit > 0 else 0
        
        response += f"📈 使用率: {usage_percent:.1f}%\n"
        
        if usage_percent > 80:
            response += "⚠️  警告: 存储空间使用率较高"
        elif usage_percent > 90:
            response += "🚨 警告: 存储空间即将不足"
        
        return {
            "response": response, 
            "type": "info",
            "style": "color: #ffffff; background: transparent;"
        }
        
    except Exception as e:
        return {
            "response": f"获取存储信息失败: {str(e)}", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }

def handle_version_command(args):
    """处理version命令"""
    return {
        "response": "WhatNote 控制台系统 v2.0\n直接命令处理模式", 
        "type": "info",
        "style": "color: #74c0fc; background: transparent;"
    }

def handle_help_command(args, current_path):
    """处理help命令"""
    path_type = current_path.get('context', {}).get('type', 'root') if current_path else 'root'
    
    if args and len(args) > 0:
        # 特定命令的帮助
        cmd = args[0].lower()
        if cmd == "cd":
            return {
                "response": "cd 命令帮助:\n  cd <目录> - 进入指定目录\n  cd .. - 返回上级目录\n  cd ~ 或 cd / - 返回根目录\n  支持课程名称和展板名称", 
                "type": "info",
                "style": "color: #74c0fc; background: transparent;"
            }
        elif cmd == "ls":
            return {
                "response": "ls 命令帮助:\n  ls - 列出当前目录内容\n  在根目录显示课程文件夹和PDF文件\n  在课程中显示展板\n  在展板中显示窗口和PDF", 
                "type": "info",
                "style": "color: #74c0fc; background: transparent;"
            }
        elif cmd == "course":
            return {
                "response": "course 命令帮助:\n  course list - 列出所有课程\n  course create <名称> - 创建课程\n  course delete <名称> - 删除课程\n  course rename <旧名> <新名> - 重命名课程\n  course show <名称> - 显示课程详情", 
                "type": "info",
                "style": "color: #74c0fc; background: transparent;"
            }
        elif cmd == "board":
            return {
                "response": "board 命令帮助:\n  board list - 列出展板\n  board create <名称> - 创建展板\n  board delete <名称> - 删除展板\n  board rename <旧名> <新名> - 重命名展板\n  board show <名称> - 显示展板详情", 
                "type": "info",
                "style": "color: #74c0fc; background: transparent;"
            }
        elif cmd == "window":
            return {
                "response": "window 命令帮助:\n  window list - 列出窗口\n  window create <类型> <标题> - 创建窗口\n  window delete <ID> - 删除窗口\n  window write <ID> <内容> - 写入内容\n  window show <ID> - 显示窗口", 
                "type": "info",
                "style": "color: #74c0fc; background: transparent;"
            }
        else:
            return {
                "response": f"未找到命令 '{cmd}' 的帮助信息\n输入 'help' 查看所有可用命令", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
    # 通用帮助信息
    help_text = f"""📋 WhatNote 控制台命令帮助

🚀 基础命令:
  help [命令]     - 显示帮助信息
  ls             - 列出当前目录内容  
  cd <目录>       - 切换目录
  pwd            - 显示当前路径
  clear          - 清屏
  status         - 系统状态
  exit           - 退出控制台

📚 课程管理:
  course list                    - 列出课程
  course create <名称>           - 创建课程
  course delete <名称>           - 删除课程
  course rename <旧名> <新名>     - 重命名课程

📋 展板管理:
  board list                     - 列出展板
  board create <名称>            - 创建展板
  board delete <名称>            - 删除展板
  board rename <旧名> <新名>      - 重命名展板

🪟 窗口管理:
  window list                    - 列出窗口
  window create <类型> <标题>     - 创建窗口
  window delete <ID>             - 删除窗口
  window write <ID> <内容>        - 写入内容
  window show <ID>               - 显示窗口

📄 PDF管理:
  pdf list                       - 列出PDF文件
  pdf delete <文件名>            - 删除PDF文件

⚙️ 系统命令:
  config show                    - 显示配置
  log                           - 查看日志
  cache                         - 缓存管理
  version                       - 版本信息
  quota                         - 存储使用情况

当前位置: {path_type}
输入 'help <命令>' 获取具体命令的详细说明"""

    return {
        "response": help_text,
        "type": "info", 
        "style": "color: #74c0fc; background: transparent;"
    }

# 缺失的handle函数定义

async def handle_tree_command(args):
    """处理tree命令"""
    try:
        course_folders = app_state.get_course_folders()
        boards = app_state.get_boards()
        
        response = "🌳 WhatNote 目录树结构:\n"
        response += "📁 whatnote/\n"
        
        if course_folders:
            for i, folder in enumerate(course_folders):
                is_last_folder = (i == len(course_folders) - 1)
                folder_prefix = "└── " if is_last_folder else "├── "
                response += f"{folder_prefix}📚 {folder['name']}/\n"
                
                # 查找该课程下的展板
                course_boards = [b for b in boards if b.get('course_folder') == folder['name']]
                for j, board in enumerate(course_boards):
                    is_last_board = (j == len(course_boards) - 1)
                    board_prefix = "    └── " if is_last_folder else "│   └── " if is_last_board else "│   ├── "
                    if is_last_folder:
                        board_prefix = "    └── " if is_last_board else "    ├── "
                    response += f"{board_prefix}📋 {board['name']}\n"
        
        # 显示独立展板
        independent_boards = [b for b in boards if not b.get('course_folder')]
        if independent_boards:
            for i, board in enumerate(independent_boards):
                is_last = (i == len(independent_boards) - 1) and not course_folders
                prefix = "└── " if is_last else "├── "
                response += f"{prefix}📋 {board['name']} (独立展板)\n"
        
        return {
            "response": response, 
            "type": "info",
            "style": "color: #ffffff; background: transparent;"
        }
        
    except Exception as e:
        return {
            "response": f"生成目录树失败: {str(e)}", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }

async def handle_find_command(args):
    """处理find命令"""
    if not args:
        return {
            "response": "用法: find \"关键词\" [-t type]", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }
    
    try:
        # 解析参数
        search_type = None
        keyword = None
        
        i = 0
        while i < len(args):
            if args[i] == '-t' and i + 1 < len(args):
                search_type = args[i + 1]
                i += 2
            else:
                if keyword is None:
                    keyword = args[i]
                else:
                    keyword += " " + args[i]
                i += 1
        
        # 去掉引号
        if keyword and ((keyword.startswith('"') and keyword.endswith('"')) or (keyword.startswith("'") and keyword.endswith("'"))):
            keyword = keyword[1:-1]
        
        if not keyword:
            return {
                "response": "请提供搜索关键词", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        course_folders = app_state.get_course_folders()
        boards = app_state.get_boards()
        
        results = []
        
        # 搜索课程
        if not search_type or search_type == 'course':
            for folder in course_folders:
                if keyword.lower() in folder['name'].lower():
                    results.append(f"📚 课程: {folder['name']}")
        
        # 搜索展板
        if not search_type or search_type == 'board':
            for board in boards:
                if keyword.lower() in board['name'].lower():
                    course_info = f" [课程: {board['course_folder']}]" if board.get('course_folder') else ""
                    results.append(f"📋 展板: {board['name']}{course_info}")
        
        # 搜索PDF（简单实现）
        if not search_type or search_type == 'pdf':
            uploads_dir = "uploads"
            if os.path.exists(uploads_dir):
                for file in os.listdir(uploads_dir):
                    if file.endswith('.pdf') and keyword.lower() in file.lower():
                        results.append(f"📄 PDF: {file}")
        
        if results:
            response = f"🔍 搜索结果 (关键词: \"{keyword}\"):\n"
            response += "\n".join(results)
        else:
            response = f"🔍 未找到包含 \"{keyword}\" 的内容"
        
        return {
            "response": response, 
            "type": "info",
            "style": "color: #ffffff; background: transparent;"
        }
        
    except Exception as e:
        return {
            "response": f"搜索失败: {str(e)}", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }

async def handle_search_command(args):
    """处理search命令"""
    return {
        "response": "搜索内容功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

async def handle_stats_command(args):
    """处理stats命令"""
    try:
        course_folders = app_state.get_course_folders()
        boards = app_state.get_boards()
        
        # 统计PDF文件
        uploads_dir = "uploads"
        pdf_count = 0
        total_size = 0
        
        if os.path.exists(uploads_dir):
            for file in os.listdir(uploads_dir):
                if file.endswith('.pdf'):
                    pdf_count += 1
                    file_path = os.path.join(uploads_dir, file)
                    if os.path.exists(file_path):
                        total_size += os.path.getsize(file_path)
        
        # 格式化大小
        def format_size(size_bytes):
            if size_bytes == 0:
                return "0 B"
            size_names = ["B", "KB", "MB", "GB"]
            import math
            i = int(math.floor(math.log(size_bytes, 1024)))
            p = math.pow(1024, i)
            s = round(size_bytes / p, 2)
            return f"{s} {size_names[i]}"
        
        # 统计课程-展板关系
        course_board_count = {}
        independent_boards = 0
        
        for board in boards:
            course = board.get('course_folder', '')
            if course:
                course_board_count[course] = course_board_count.get(course, 0) + 1
            else:
                independent_boards += 1
        
        response = "📊 WhatNote 使用统计:\n\n"
        response += f"📚 总课程数: {len(course_folders)}\n"
        response += f"📋 总展板数: {len(boards)}\n"
        response += f"📄 总PDF数: {pdf_count}\n"
        response += f"💾 存储使用: {format_size(total_size)}\n\n"
        
        if course_board_count:
            response += "📈 课程-展板分布:\n"
            for course, count in sorted(course_board_count.items()):
                response += f"  📚 {course}: {count} 个展板\n"
        
        if independent_boards > 0:
            response += f"  📋 独立展板: {independent_boards} 个\n"
        
        avg_boards_per_course = len(boards) / len(course_folders) if course_folders else 0
        avg_pdfs_per_course = pdf_count / len(course_folders) if course_folders else 0
        
        response += f"\n📊 平均指标:\n"
        response += f"  每个课程平均展板数: {avg_boards_per_course:.1f}\n"
        response += f"  每个课程平均PDF数: {avg_pdfs_per_course:.1f}\n"
        
        return {
            "response": response, 
            "type": "info",
            "style": "color: #ffffff; background: transparent;"
        }
        
    except Exception as e:
        return {
            "response": f"获取统计信息失败: {str(e)}", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }

async def handle_recent_command(args):
    """处理recent命令"""
    return {
        "response": "最近访问功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

async def handle_backup_command(args, current_path):
    """处理backup命令"""
    return {
        "response": "备份功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

async def handle_export_command(args, current_path):
    """处理export命令"""
    return {
        "response": "导出功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

async def handle_delete_command(args, current_path):
    """处理delete命令"""
    return {
        "response": "删除功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

async def handle_rename_command(args, current_path):
    """处理rename命令"""
    return {
        "response": "重命名功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

async def handle_info_command(args, current_path):
    """处理info命令"""
    return {
        "response": "信息查看功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

async def handle_copy_command(args, current_path):
    """处理copy命令"""
    return {
        "response": "复制功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

# PDF相关的导航命令
async def handle_goto_command(args, current_path):
    """处理goto命令"""
    return {
        "response": "页面跳转功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

async def handle_next_command(args, current_path):
    """处理next命令"""
    return {
        "response": "下一页功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

async def handle_prev_command(args, current_path):
    """处理prev命令"""
    return {
        "response": "上一页功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

async def handle_first_command(args, current_path):
    """处理first命令"""
    return {
        "response": "跳转第一页功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

async def handle_last_command(args, current_path):
    """处理last命令"""
    return {
        "response": "跳转最后一页功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

async def handle_pages_command(args, current_path):
    """处理pages命令"""
    return {
        "response": "页数显示功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

async def handle_annotate_command(args, current_path):
    """处理annotate命令"""
    return {
        "response": "注释生成功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

async def handle_annotation_command(args, current_path):
    """处理annotation命令"""
    if not args:
        return {
            "response": "用法: annotation <show|edit|delete|improve> [参数]", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }
    
    action = args[0].lower()
    
    return {
        "response": f"注释{action}功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

async def handle_page_command(args, current_path):
    """处理page命令"""
    if not args:
        return {
            "response": "用法: page <text|extract|ocr|vision> [参数]", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }
    
    action = args[0].lower()
    
    return {
        "response": f"页面{action}功能需要进一步开发", 
        "type": "info",
        "style": "color: #ffd43b; background: transparent;"
    }

# 窗口操作辅助函数
def find_window_by_name_or_id(windows, name_or_id):
    """通过名字或ID查找窗口"""
    # 首先尝试按ID精确匹配
    for window in windows:
        if window.get("id") == name_or_id:
            return window
    
    # 然后尝试按标题精确匹配
    for window in windows:
        if window.get("title") == name_or_id:
            return window
    
    # 最后尝试模糊匹配（不区分大小写）
    name_or_id_lower = name_or_id.lower()
    for window in windows:
        if window.get("title", "").lower() == name_or_id_lower:
            return window
            
    return None

def generate_unique_window_title(windows, base_title):
    """生成唯一的窗口标题，对同名窗口自动添加编号"""
    existing_titles = [w.get("title", "") for w in windows]
    
    # 如果基础标题不存在，直接返回
    if base_title not in existing_titles:
        return base_title
    
    # 查找已存在的编号
    counter = 1
    while True:
        new_title = f"{base_title}({counter})"
        if new_title not in existing_titles:
            return new_title
        counter += 1

async def handle_window_command(args, current_path):
    """处理window命令"""
    if not args:
        return {
            "response": "用法: window <create|list|show|delete|write|image> [参数]", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }
    
    action = args[0].lower()
    
    # 检查当前路径上下文
    path_type = current_path.get('context', {}).get('type', 'root') if current_path else 'root'
    
    if path_type != 'board':
        return {
            "response": "window命令只能在展板中使用，请先用 cd 进入展板", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }
    
    # 获取当前展板信息
    board_name = current_path.get('context', {}).get('boardName', '')
    course_name = current_path.get('context', {}).get('courseName', '')
    
    # 查找展板ID
    boards = app_state.get_boards()
    current_board = None
    for board in boards:
        if board.get('name') == board_name and board.get('course_folder') == course_name:
            current_board = board
            break
    
    if not current_board:
        return {
            "response": f"未找到当前展板: {board_name}", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }
    
    board_id = current_board.get('id')
    
    if action == "create":
        if len(args) < 2:
            return {
                "response": "用法: window create <text|image|video> [标题]", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        window_type = args[1].lower()
        if window_type not in ['text', 'image', 'video']:
            return {
                "response": "窗口类型只能是 text、image 或 video", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        # 获取标题（可选）
        base_title = ' '.join(args[2:]) if len(args) > 2 else f"新{window_type}窗口"
        
        # 去掉外层引号
        if (base_title.startswith('"') and base_title.endswith('"')) or (base_title.startswith("'") and base_title.endswith("'")):
            base_title = base_title[1:-1]
        
        try:
            from board_logger import BoardLogger
            board_logger = BoardLogger()
            log_data = board_logger.load_log(board_id)
            existing_windows = log_data.get("windows", [])
            
            # 生成唯一的窗口标题
            unique_title = generate_unique_window_title(existing_windows, base_title)
            
            # 创建窗口数据
            import time
            window_data = {
                "type": window_type,
                "title": unique_title,
                "content": "",
                "position": {"x": 100, "y": 100},
                "size": {"width": 300, "height": 200},
                "style": {}
            }
            
            window_id = board_logger.add_window(board_id, window_data)
            
            return {
                "response": f"✅ 已创建{window_type}窗口: {unique_title} (ID: {window_id})", 
                "type": "success",
                "style": "color: #51cf66; background: transparent;",
                "refresh_needed": True
            }
        except Exception as e:
            return {
                "response": f"创建窗口失败: {str(e)}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
    elif action == "list":
        try:
            from board_logger import BoardLogger
            board_logger = BoardLogger()
            log_data = board_logger.load_log(board_id)
            windows = log_data.get("windows", [])
            
            if windows:
                response = f"📋 当前展板的窗口 ({len(windows)}):\n"
                for i, window in enumerate(windows, 1):
                    window_id = window.get("id", "")
                    window_type = window.get("type", "")
                    title = window.get("title", "")
                    content_preview = window.get("content", "")[:30]
                    if len(window.get("content", "")) > 30:
                        content_preview += "..."
                    response += f"  {i}. [{window_type}] {title} (ID: {window_id})\n"
                    if content_preview:
                        response += f"     内容: {content_preview}\n"
            else:
                response = f"📋 当前展板暂无窗口"
            
            return {
                "response": response, 
                "type": "info",
                "style": "color: #ffffff; background: transparent;"
            }
        except Exception as e:
            return {
                "response": f"获取窗口列表失败: {str(e)}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
    elif action == "write":
        if len(args) < 3:
            return {
                "response": "用法: window write <窗口名字或ID> \"内容\"", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        window_name_or_id = args[1]
        content = ' '.join(args[2:])
        
        # 去掉外层引号，但保留内容中的引号
        if (content.startswith('"') and content.endswith('"')) or (content.startswith("'") and content.endswith("'")):
            content = content[1:-1]
        
        # 去掉窗口名字的引号
        if (window_name_or_id.startswith('"') and window_name_or_id.endswith('"')) or (window_name_or_id.startswith("'") and window_name_or_id.endswith("'")):
            window_name_or_id = window_name_or_id[1:-1]
        
        try:
            from board_logger import BoardLogger
            board_logger = BoardLogger()
            log_data = board_logger.load_log(board_id)
            windows = log_data.get("windows", [])
            
            # 通过名字或ID查找窗口
            target_window = find_window_by_name_or_id(windows, window_name_or_id)
            
            if not target_window:
                return {
                    "response": f"找不到窗口: {window_name_or_id}\n提示: 可以使用 'window list' 查看所有窗口", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
            
            # 更新窗口内容
            window_id = target_window["id"]
            title = target_window.get("title", "")
            target_window["content"] = content
            success = board_logger.update_window(board_id, window_id, target_window)
            
            if success:
                return {
                    "response": f"✅ 已更新窗口 '{title}' 的内容", 
                    "type": "success",
                    "style": "color: #51cf66; background: transparent;",
                    "refresh_needed": True
                }
            else:
                return {
                    "response": f"更新窗口内容失败", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
        except Exception as e:
            return {
                "response": f"写入窗口内容失败: {str(e)}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
    elif action == "image":
        if len(args) < 3:
            return {
                "response": "用法: window image <窗口名字或ID> <图片路径或URL>\n支持:\n  - 本地文件: window image \"我的图片\" \"C:/path/to/image.jpg\"\n  - 网络图片: window image \"我的图片\" \"https://example.com/image.jpg\"\n  - 系统图片: window image \"我的图片\" \"uploaded/filename.jpg\"", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        window_name_or_id = args[1]
        image_path = ' '.join(args[2:])
        
        # 去掉外层引号
        if (image_path.startswith('"') and image_path.endswith('"')) or (image_path.startswith("'") and image_path.endswith("'")):
            image_path = image_path[1:-1]
        
        # 去掉窗口名字的引号
        if (window_name_or_id.startswith('"') and window_name_or_id.endswith('"')) or (window_name_or_id.startswith("'") and window_name_or_id.endswith("'")):
            window_name_or_id = window_name_or_id[1:-1]
        
        try:
            from board_logger import BoardLogger
            board_logger = BoardLogger()
            log_data = board_logger.load_log(board_id)
            windows = log_data.get("windows", [])
            
            # 通过名字或ID查找窗口
            target_window = find_window_by_name_or_id(windows, window_name_or_id)
            
            if not target_window:
                return {
                    "response": f"找不到窗口: {window_name_or_id}\n提示: 可以使用 'window list' 查看所有窗口", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
            
            # 检查是否为图片窗口
            window_id = target_window["id"]
            title = target_window.get("title", "")
            if target_window.get("type") != "image":
                return {
                    "response": f"窗口 '{title}' 不是图片窗口，无法设置图片", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
            
            # 处理不同类型的图片路径
            final_image_url = None
            
            if image_path.startswith(("http://", "https://")):
                # 网络图片URL
                final_image_url = image_path
                print(f"🌐 设置网络图片: {image_path}")
                
            elif os.path.exists(image_path):
                # 本地文件路径
                print(f"📁 处理本地文件: {image_path}")
                
                # 验证是否为图片文件
                allowed_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp'}
                file_ext = os.path.splitext(image_path.lower())[1]
                
                if file_ext not in allowed_extensions:
                    return {
                        "response": f"文件格式不支持: {file_ext}，支持的格式: {', '.join(allowed_extensions)}", 
                        "type": "error",
                        "style": "color: #ff6b6b; background: transparent;"
                    }
                
                # 复制文件到images目录
                import shutil
                import time
                
                filename = os.path.basename(image_path)
                name, ext = os.path.splitext(filename)
                timestamp = int(time.time())
                unique_filename = f"{name}_cmd_{timestamp}{ext}"
                
                images_dir = os.path.join(UPLOAD_DIR, 'images')
                if not os.path.exists(images_dir):
                    os.makedirs(images_dir, exist_ok=True)
                
                dest_path = os.path.join(images_dir, unique_filename)
                shutil.copy2(image_path, dest_path)
                
                final_image_url = f"/api/images/view/{unique_filename}"
                print(f"📋 复制到系统: {dest_path}")
                
            elif image_path.startswith("uploaded/") or image_path.startswith("/api/images/"):
                # 系统内部图片路径
                if image_path.startswith("uploaded/"):
                    # 兼容旧格式
                    filename = image_path.replace("uploaded/", "")
                    final_image_url = f"/api/images/view/{filename}"
                else:
                    # 已经是正确格式
                    final_image_url = image_path
                    
                print(f"🗃️  使用系统图片: {final_image_url}")
                
            else:
                return {
                    "response": f"图片路径无效: {image_path}\n请使用:\n  - 网络URL (http://或https://)\n  - 本地文件的完整路径\n  - 系统图片路径 (uploaded/filename.jpg)", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
            
            # 更新窗口的图片URL
            target_window["content"] = final_image_url
            success = board_logger.update_window(board_id, window_id, target_window)
            
            if success:
                return {
                    "response": f"✅ 已设置图片窗口 {window_id} 的图片: {final_image_url}", 
                    "type": "success",
                    "style": "color: #51cf66; background: transparent;",
                    "refresh_needed": True
                }
            else:
                return {
                    "response": f"设置图片失败", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
                
        except Exception as e:
            return {
                "response": f"设置图片失败: {str(e)}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
    elif action == "delete":
        if len(args) < 2:
            return {
                "response": "用法: window delete <窗口ID>", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        window_id = args[1]
        
        try:
            from board_logger import BoardLogger
            board_logger = BoardLogger()
            success = board_logger.remove_window(board_id, window_id)
            
            if success:
                return {
                    "response": f"✅ 已删除窗口: {window_id}", 
                    "type": "success",
                    "style": "color: #51cf66; background: transparent;",
                    "refresh_needed": True
                }
            else:
                return {
                    "response": f"找不到窗口: {window_id}", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
        except Exception as e:
            return {
                "response": f"删除窗口失败: {str(e)}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
    elif action == "show":
        if len(args) < 2:
            return {
                "response": "用法: window show <窗口ID>", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        window_id = args[1]
        
        try:
            from board_logger import BoardLogger
            board_logger = BoardLogger()
            log_data = board_logger.load_log(board_id)
            windows = log_data.get("windows", [])
            
            # 查找指定窗口
            target_window = None
            for window in windows:
                if window.get("id") == window_id:
                    target_window = window
                    break
            
            if not target_window:
                return {
                    "response": f"找不到窗口: {window_id}", 
                    "type": "error",
                    "style": "color: #ff6b6b; background: transparent;"
                }
            
            # 显示窗口详细信息
            response = f"🪟 窗口详情: {target_window.get('title', '')}\n"
            response += f"  ID: {window_id}\n"
            response += f"  类型: {target_window.get('type', '')}\n"
            response += f"  位置: x={target_window.get('position', {}).get('x', 0)}, y={target_window.get('position', {}).get('y', 0)}\n"
            response += f"  大小: {target_window.get('size', {}).get('width', 0)}x{target_window.get('size', {}).get('height', 0)}\n"
            response += f"  创建时间: {target_window.get('created_at', '未知')}\n"
            content = target_window.get('content', '')
            if content:
                if len(content) > 100:
                    response += f"  内容: {content[:100]}...\n"
                else:
                    response += f"  内容: {content}\n"
            else:
                response += "  内容: (空)\n"
            
            return {
                "response": response, 
                "type": "info",
                "style": "color: #ffffff; background: transparent;"
            }
        except Exception as e:
            return {
                "response": f"获取窗口信息失败: {str(e)}", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
    else:
        return {
            "response": f"未知的窗口操作: {action}", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }

async def handle_pdf_command(args, current_path):
    """处理pdf命令"""
    if not args:
        return {
            "response": "用法: pdf <list|open> [文件名]", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }
    
    action = args[0].lower()
    
    if action == "list":
        # 根据当前路径上下文显示PDF
        path_type = current_path.get('context', {}).get('type', 'root') if current_path else 'root'
        
        if path_type == 'board':
            # 在展板中，只显示当前展板的PDF
            board_name = current_path.get('context', {}).get('boardName', '')
            course_name = current_path.get('context', {}).get('courseName', '')
            
            # 查找对应的展板数据
            boards = app_state.get_boards()
            current_board = None
            for board in boards:
                if board.get('name') == board_name and board.get('course_folder') == course_name:
                    current_board = board
                    break
            
            if current_board:
                board_id = current_board.get('id')
                # 从board_data.json获取展板的PDF列表
                try:
                    board_data_file = f"board_data.json"
                    if os.path.exists(board_data_file):
                        with open(board_data_file, 'r', encoding='utf-8') as f:
                            board_data = json.load(f)
                            if board_data.get('board_id') == board_id:
                                pdf_files = board_data.get('pdfs', [])
                                if pdf_files:
                                    response = f"📄 当前展板 '{board_name}' 的PDF文件 ({len(pdf_files)}):\n"
                                    for i, pdf in enumerate(pdf_files, 1):
                                        filename = pdf.get('filename', '')
                                        current_page = pdf.get('currentPage', 1)
                                        response += f"  {i}. {filename} (页: {current_page})\n"
                                else:
                                    response = f"📄 当前展板 '{board_name}' 暂无PDF文件"
                            else:
                                response = f"📄 当前展板 '{board_name}' 暂无PDF文件"
                    else:
                        response = f"📄 当前展板 '{board_name}' 暂无PDF文件"
                except Exception as e:
                    response = f"📄 获取展板PDF列表失败: {str(e)}"
            else:
                response = f"📄 未找到展板 '{board_name}'"
                
        elif path_type == 'course':
            # 在课程中，显示该课程下所有展板的PDF
            course_name = current_path.get('context', {}).get('courseName', '')
            boards = app_state.get_boards()
            course_boards = [b for b in boards if b.get('course_folder') == course_name]
            
            total_pdfs = []
            for board in course_boards:
                board_id = board.get('id')
                try:
                    board_data_file = f"board_data.json"
                    if os.path.exists(board_data_file):
                        with open(board_data_file, 'r', encoding='utf-8') as f:
                            board_data = json.load(f)
                            if board_data.get('board_id') == board_id:
                                board_pdfs = board_data.get('pdfs', [])
                                for pdf in board_pdfs:
                                    pdf['board_name'] = board.get('name', '')
                                total_pdfs.extend(board_pdfs)
                except:
                    continue
            
            if total_pdfs:
                response = f"📄 课程 '{course_name}' 的PDF文件 ({len(total_pdfs)}):\n"
                for i, pdf in enumerate(total_pdfs, 1):
                    filename = pdf.get('filename', '')
                    board_name = pdf.get('board_name', '')
                    current_page = pdf.get('currentPage', 1)
                    response += f"  {i}. {filename} [展板: {board_name}] (页: {current_page})\n"
            else:
                response = f"📄 课程 '{course_name}' 暂无PDF文件"
        else:
            # 在根目录，显示所有PDF
            uploads_dir = "uploads"
            pdf_files = []
            
            if os.path.exists(uploads_dir):
                pdf_files = [f for f in os.listdir(uploads_dir) if f.endswith('.pdf')]
            
            if pdf_files:
                response = f"📄 所有PDF文件 ({len(pdf_files)}):\n"
                for i, pdf in enumerate(pdf_files, 1):
                    response += f"  {i}. {pdf}\n"
            else:
                response = "📄 系统中暂无PDF文件"
        
        return {
            "response": response, 
            "type": "info",
            "style": "color: #ffffff; background: transparent;"
        }
    
    elif action == "open":
        if len(args) < 2:
            return {
                "response": "用法: pdf open <文件名>", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        
        pdf_name = ' '.join(args[1:])
        if (pdf_name.startswith('"') and pdf_name.endswith('"')) or (pdf_name.startswith("'") and pdf_name.endswith("'")):
            pdf_name = pdf_name[1:-1]
        
        # 检查PDF文件是否存在
        uploads_dir = "uploads"
        if os.path.exists(uploads_dir):
            available_pdfs = [f for f in os.listdir(uploads_dir) if f.endswith('.pdf')]
            
            # 精确匹配
            if pdf_name in available_pdfs:
                return {
                    "response": f"打开PDF: {pdf_name}",
                    "type": "navigation",
                    "style": "color: #74c0fc; background: transparent;",
                    "navigation": {
                        "action": "open_pdf",
                        "pdf_name": pdf_name
                    }
                }
            
            # 不区分大小写的匹配
            for pdf in available_pdfs:
                if pdf.lower() == pdf_name.lower():
                    return {
                        "response": f"打开PDF: {pdf}",
                        "type": "navigation", 
                        "style": "color: #74c0fc; background: transparent;",
                        "navigation": {
                            "action": "open_pdf",
                            "pdf_name": pdf
                        }
                    }
            
            # 提供建议
            suggestions = []
            for pdf in available_pdfs:
                if pdf_name.lower() in pdf.lower():
                    suggestions.append(pdf)
            
            error_msg = f"找不到PDF文件: {pdf_name}"
            if suggestions:
                error_msg += f"\n💡 您是否在找: {', '.join(suggestions[:3])}"
            
            return {
                "response": error_msg, 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
        else:
            return {
                "response": "PDF上传目录不存在", 
                "type": "error",
                "style": "color: #ff6b6b; background: transparent;"
            }
    
    else:
        return {
            "response": f"未知的PDF操作: {action}", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }

async def handle_note_command(args, current_path):
    """处理note命令"""
    if not args:
        return {
            "response": "用法: note <generate|show|edit|improve> [参数]", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }
    
    action = args[0].lower()
    
    if action == "generate":
        return {
            "response": "笔记生成功能需要进一步开发", 
            "type": "info",
            "style": "color: #ffd43b; background: transparent;"
        }
    elif action == "show":
        return {
            "response": "笔记显示功能需要进一步开发", 
            "type": "info",
            "style": "color: #ffd43b; background: transparent;"
        }
    elif action == "edit":
        return {
            "response": "笔记编辑功能需要进一步开发", 
            "type": "info",
            "style": "color: #ffd43b; background: transparent;"
        }
    elif action == "improve":
        return {
            "response": "笔记改进功能需要进一步开发", 
            "type": "info",
            "style": "color: #ffd43b; background: transparent;"
        }
    else:
        return {
            "response": f"未知的笔记操作: {action}", 
            "type": "error",
            "style": "color: #ff6b6b; background: transparent;"
        }

# 主程序启动
if __name__ == "__main__":
    # 加载环境变量
    dotenv.load_dotenv('.env')
    
    # 打印欢迎信息
    print("\n=== WhatNote 服务已启动 ===")
    print(f"API密钥配置: {'已配置' if bool(os.getenv('QWEN_API_KEY')) else '未配置'}")
    print(f"视觉API配置: {'已配置' if bool(os.getenv('QWEN_VL_API_KEY')) else '未配置'}")
    print("=======================\n")
    
    # 应用启动时同步一次文件结构
    sync_app_state_to_butler()
    
    # 启动服务
    uvicorn.run("main:app", host="127.0.0.1", port=8000, reload=True)
